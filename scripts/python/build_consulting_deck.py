from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
from pptx.oxml.ns import qn
import os

KSP = r"E:\N8N\scraper\scraper2\csi-ultimate\ksp_images"
CHART = r"E:\N8N\scraper\scraper2\csi-ultimate\consulting_charts"
OUTPUT = r"E:\N8N\scraper\scraper2\csi-ultimate\Construction_Failures_Executive_Deck.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

NAVY = RGBColor(27, 42, 74)
DARK = RGBColor(15, 25, 50)
GOLD = RGBColor(200, 150, 46)
TEAL = RGBColor(26, 138, 158)
WHITE = RGBColor(255, 255, 255)
LGRAY = RGBColor(200, 200, 210)
MGRAY = RGBColor(150, 150, 165)
RED = RGBColor(192, 57, 43)
GREEN = RGBColor(39, 174, 96)
ORANGE = RGBColor(230, 126, 34)
SOFT = RGBColor(240, 240, 245)
DRED = RGBColor(142, 27, 27)
BG = RGBColor(248, 249, 251)

FONT = 'Calibri'
FLT = 'Calibri Light'

def set_alpha(shape, val):
    spPr = shape._element.spPr
    sf = spPr.find(qn('a:solidFill'))
    if sf is not None:
        srgb = sf.find(qn('a:srgbClr'))
        if srgb is not None:
            srgb.set('alpha', str(int(val * 1000)))

def R(l, t, w, h, c, a=None):
    s = prs.slides[-1].shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()
    if a: set_alpha(s, a)
    return s

def T(l, t, w, h, txt, sz=14, c=WHITE, b=False, al=PP_ALIGN.LEFT, fn=FONT):
    tb = prs.slides[-1].shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = txt; p.font.size = Pt(sz)
    p.font.color.rgb = c; p.font.bold = b; p.font.name = fn; p.alignment = al
    return tb

def TX(l, t, w, h, lines, al=PP_ALIGN.LEFT):
    tb = prs.slides[-1].shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, (txt, sz, c, bld, fn) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt; p.font.size = Pt(sz); p.font.color.rgb = c
        p.font.bold = bld; p.font.name = fn; p.alignment = al; p.space_after = Pt(3)
    return tb

def IMG(path, l=0, t=0, w=None, h=None):
    if w is None: w = SW
    if h is None: h = SH
    prs.slides[-1].shapes.add_picture(path, l, t, w, h)

def C(p):
    return os.path.join(CHART, p)

def K(p):
    return os.path.join(KSP, p)

def SL():
    return prs.slides.add_slide(prs.slide_layouts[6])

# ══════════════════════════════════════════════════════════════
# SLIDE 1  —  EXECUTIVE COVER
# ══════════════════════════════════════════════════════════════
SL()
IMG(K('1509_KSP_Park_Overview_001_1920x1080.jpg'))
R(0, 0, SW, SH, DARK, 0.55)
R(0, 0, Inches(0.15), SH, GOLD)
# Gold band
R(Inches(1.2), Inches(1.8), Inches(2), Pt(5), GOLD)
TX(Inches(1.2), Inches(2.1), Inches(11), Inches(3.5), [
    ('THE HIDDEN CAUSES BEHIND', 14, LGRAY, False, FLT),
    ('CONSTRUCTION PROJECT FAILURES', 42, WHITE, True, FLT),
    ('', 10, WHITE, False, FLT),
    ('Why 70% of Mega Projects Exceed Budget', 20, GOLD, False, FLT),
    ('and 60% Miss Their Deadlines', 20, GOLD, False, FLT),
], PP_ALIGN.LEFT)
T(Inches(1.2), Inches(5.5), Inches(10), Inches(0.5),
  'Executive Brief  |  PMO Advisory  |  Project Controls Intelligence', 12, MGRAY, False, PP_ALIGN.LEFT, FLT)
T(Inches(1.2), Inches(5.9), Inches(10), Inches(0.5),
  'Confidential  |  Board-Level Review', 11, MGRAY, False, PP_ALIGN.LEFT, FLT)
# Bottom-right accent
R(Inches(10.5), Inches(6.2), Inches(2.5), Pt(2), GOLD)

# ══════════════════════════════════════════════════════════════
# SLIDE 2  —  EXECUTIVE SNAPSHOT (infographic cards)
# ══════════════════════════════════════════════════════════════
SL()
R(0, 0, SW, SH, WHITE)
# Top bar
R(0, 0, SW, Inches(0.08), NAVY)
R(0, Inches(0.08), SW, Pt(3), GOLD)

T(Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
  'The Construction Industry Crisis  |  By the Numbers', 22, NAVY, True, PP_ALIGN.LEFT, FLT)

# 4 big KPI cards
cards = [
    ('70%', 'of projects exceed\ntheir approved budget', RED, 'Cost Overrun\nEpidemic', '+38% average\noverrun'),
    ('60%', 'miss their original\ncompletion deadlines', RED, 'Schedule\nFailure Rate', '+42% average\ndelay'),
    ('85%', 'of executives lack\nearly warning systems', GOLD, 'Visibility\nGap', 'No real-time\nEVM data'),
    ('$1.2T', 'wasted annually on\ndelayed construction', DRED, 'Global\nImpact', 'Across all\nsectors'),
]
for i, (big, desc, color, title, sub) in enumerate(cards):
    x = Inches(0.5) + Inches(i * 3.2)
    # Card background
    R(x, Inches(1.5), Inches(2.9), Inches(3.8), WHITE)
    R(x, Inches(1.5), Inches(2.9), Inches(0.08), color)
    # Big number
    T(x + Inches(0.2), Inches(1.8), Inches(2.5), Inches(0.9), big, 48, color, True, PP_ALIGN.LEFT, FLT)
    # Description
    T(x + Inches(0.2), Inches(2.7), Inches(2.5), Inches(0.9), desc, 12, RGBColor(80, 80, 80), False, PP_ALIGN.LEFT, FLT)
    # Divider
    R(x + Inches(0.2), Inches(3.7), Inches(2.5), Pt(1), RGBColor(220, 220, 220))
    # Title
    T(x + Inches(0.2), Inches(3.9), Inches(2.5), Inches(0.4), title, 10, color, True, PP_ALIGN.LEFT, FLT)
    # Sub
    T(x + Inches(0.2), Inches(4.3), Inches(2.5), Inches(0.4), sub, 10, MGRAY, False, PP_ALIGN.LEFT, FLT)

# Bottom insight bar
R(Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.2), NAVY)
T(Inches(0.8), Inches(5.9), Inches(11.7), Inches(0.4),
  '"The construction industry has a productivity problem that costs the global economy $1.6 trillion per year."',
  13, GOLD, False, PP_ALIGN.LEFT, FLT)
T(Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.4),
  '— McKinsey Global Institute, 2023  |  These failures are predictable, preventable, and manageable.',
  11, LGRAY, False, PP_ALIGN.LEFT, FLT)

# ══════════════════════════════════════════════════════════════
# SLIDE 3  —  THE PROBLEM (triangle infographic)
# ══════════════════════════════════════════════════════════════
SL()
R(0, 0, SW, SH, RGBColor(245, 246, 248))
R(0, 0, Inches(0.08), SH, GOLD)

T(Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
  'The Triple Constraint Collapse', 24, NAVY, True, PP_ALIGN.LEFT, FLT)
T(Inches(0.8), Inches(0.9), Inches(11), Inches(0.4),
  'When budget, schedule, and quality deteriorate simultaneously', 13, MGRAY, False, PP_ALIGN.LEFT, FLT)

# Three connected boxes forming a triangle layout
# Left: Budget
R(Inches(0.8), Inches(4.5), Inches(3.5), Inches(1.8), RED)
T(Inches(1.0), Inches(4.6), Inches(3.1), Inches(0.4), 'BUDGET', 14, WHITE, True, PP_ALIGN.CENTER)
T(Inches(1.0), Inches(5.0), Inches(3.1), Inches(0.5), 'Budget overruns\nare the #1 cause\nof project failure.', 12, WHITE, False, PP_ALIGN.CENTER, FLT)

# Right: Schedule
R(Inches(9.0), Inches(4.5), Inches(3.5), Inches(1.8), GOLD)
T(Inches(9.2), Inches(4.6), Inches(3.1), Inches(0.4), 'SCHEDULE', 14, WHITE, True, PP_ALIGN.CENTER)
T(Inches(9.2), Inches(5.0), Inches(3.1), Inches(0.5), 'Schedule delays\ncompound costs\nand erode trust.', 12, WHITE, False, PP_ALIGN.CENTER, FLT)

# Bottom center: Quality
R(Inches(4.9), Inches(5.0), Inches(3.5), Inches(1.3), TEAL)
T(Inches(5.1), Inches(5.1), Inches(3.1), Inches(0.4), 'QUALITY', 14, WHITE, True, PP_ALIGN.CENTER)
T(Inches(5.1), Inches(5.4), Inches(3.1), Inches(0.5), 'Quality suffers when\nboth time and money\nare exhausted.', 12, WHITE, False, PP_ALIGN.CENTER, FLT)



# Center circle
R(Inches(5.0), Inches(1.6), Inches(3.3), Inches(3.2), NAVY)
T(Inches(5.2), Inches(1.8), Inches(2.9), Inches(0.5), 'THE PROBLEM', 12, GOLD, True, PP_ALIGN.CENTER)
T(Inches(5.2), Inches(2.2), Inches(2.9), Inches(0.4), 'Triple Constraint\nCollapse', 18, WHITE, True, PP_ALIGN.CENTER, FLT)
T(Inches(5.2), Inches(3.0), Inches(2.9), Inches(1.5),
  'When budget pressure\nforces schedule\ncompression,\nquality is sacrificed\ncreating a cycle\nof rework and\nescalating costs.',
  11, LGRAY, False, PP_ALIGN.CENTER, FLT)

# ══════════════════════════════════════════════════════════════
# SLIDE 4  —  WHAT EXECUTIVES USUALLY SEE (fake dashboard)
# ══════════════════════════════════════════════════════════════
SL()
R(0, 0, SW, SH, DARK)
R(0, 0, Inches(0.08), SH, GOLD)

T(Inches(0.8), Inches(0.3), Inches(8), Inches(0.5),
  'What Executives Usually See  |  The Dashboard Mirage', 24, WHITE, True, PP_ALIGN.LEFT, FLT)
T(Inches(0.8), Inches(0.75), Inches(8), Inches(0.4),
  'Green indicators that hide underlying project distress', 12, LGRAY, False, PP_ALIGN.LEFT, FLT)

# Fake dashboard image
IMG(C('fake_dashboard.png'), Inches(0.5), Inches(1.3), Inches(12.3), Inches(4.8))

# Reality check overlay bar
R(Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.9), RED)
T(Inches(0.8), Inches(6.35), Inches(11.7), Inches(0.4),
  'REALITY:  The project is already 8% over budget and 12% behind schedule  |  Dashboard shows "green"',
  16, WHITE, True, PP_ALIGN.CENTER)
T(Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.3),
  'This illusion of control is the most dangerous stage of project failure — because no corrective action is taken.',
  11, RGBColor(255, 200, 200), False, PP_ALIGN.CENTER, FLT)

# ══════════════════════════════════════════════════════════════
# SLIDE 5  —  WHAT IS ACTUALLY HAPPENING (timeline)
# ══════════════════════════════════════════════════════════════
SL()
R(0, 0, SW, SH, WHITE)
R(0, 0, Inches(0.08), SH, GOLD)

T(Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
  'What Is Actually Happening  |  Delay Accumulation Timeline', 22, NAVY, True, PP_ALIGN.LEFT, FLT)
T(Inches(0.8), Inches(0.75), Inches(11), Inches(0.4),
  'How small delays compound into catastrophic failure', 12, MGRAY, False, PP_ALIGN.LEFT, FLT)

# Main timeline bar
R(Inches(1.5), Inches(2.8), Inches(10.5), Pt(4), GOLD)

# Timeline nodes
nodes = [
    ('Month 3', 'Design\nRevision', '2 weeks\ndelay', GOLD),
    ('Month 6', 'Permit\nDelay', '4 weeks\ndelay', ORANGE),
    ('Month 9', 'Material\nShortage', '6 weeks\ndelay', RED),
    ('Month 12', 'Contractor\nDefault', '8 weeks\ndelay', RED),
    ('Month 18', 'Scope\nCreep', '10 weeks\ndelay', DRED),
    ('Month 24', 'Project\nDelivery', '30 weeks\ntotal delay', DRED),
]

for i, (month, event, impact, color) in enumerate(nodes):
    x = Inches(1.2) + Inches(i * 1.9)
    # Node dot
    shape = prs.slides[-1].shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.6), Inches(2.6), Inches(0.35), Inches(0.35))
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background()
    # Month above
    T(x + Inches(0.2), Inches(2.05), Inches(1.3), Inches(0.3), month, 9, NAVY, True, PP_ALIGN.CENTER)
    # Event above that
    T(x + Inches(0.05), Inches(1.45), Inches(1.5), Inches(0.5), event, 11, RGBColor(60, 60, 60), True, PP_ALIGN.CENTER, FLT)
    # Impact below
    T(x + Inches(0.05), Inches(3.1), Inches(1.5), Inches(0.4), impact, 10, color, True, PP_ALIGN.CENTER, FLT)

# Delay accumulation bar
R(Inches(1.5), Inches(4.0), Inches(10.5), Inches(0.6), RGBColor(245, 245, 250))
# Stacked delay visualization
delays_pct = [0.05, 0.15, 0.15, 0.20, 0.25, 0.20]  # relative contribution
x_start = Inches(1.5)
colors_delay = [GOLD, ORANGE, RED, RED, DRED, DRED]
for i, (d, c) in enumerate(zip(delays_pct, colors_delay)):
    w = Inches(10.5 * d)
    R(x_start, Inches(4.0), w, Inches(0.6), c, 0.85)
    x_start += w

T(Inches(1.5), Inches(4.85), Inches(10.5), Inches(0.4),
  'Cumulative delay: 30 weeks  |  Each delay compounds on the previous  |  Recovery becomes exponentially harder',
  11, NAVY, True, PP_ALIGN.CENTER, FLT)

# Bottom insight
R(Inches(1.5), Inches(5.6), Inches(10.5), Inches(1.3), NAVY)
T(Inches(1.8), Inches(5.7), Inches(10), Inches(0.5),
  'Key Insight: Each week of delay in the first 6 months costs 3 weeks at the end',
  14, GOLD, True, PP_ALIGN.LEFT)
T(Inches(1.8), Inches(6.1), Inches(10), Inches(0.5),
  'Early detection is the only way to prevent catastrophic schedule failure. Most organizations detect delays.\nonly when recovery is no longer possible — typically at 60-70% project completion.',
  11, LGRAY, False, PP_ALIGN.LEFT, FLT)

# ══════════════════════════════════════════════════════════════
# SLIDE 6  —  ROOT CAUSE ANALYSIS (fishbone diagram)
# ══════════════════════════════════════════════════════════════
SL()
R(0, 0, SW, SH, WHITE)
R(0, 0, Inches(0.08), SH, GOLD)

T(Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
  'Root Cause Analysis  |  Why Projects Fail', 22, NAVY, True, PP_ALIGN.LEFT, FLT)
T(Inches(0.8), Inches(0.75), Inches(11), Inches(0.4),
  'Ishikawa fishbone analysis of the six primary failure drivers', 12, MGRAY, False, PP_ALIGN.LEFT, FLT)

# Main spine
R(Inches(1.5), Inches(3.9), Inches(10), Pt(3), NAVY)
# Arrow head
arr = prs.slides[-1].shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(11.2), Inches(3.65), Inches(0.4), Inches(0.5))
arr.fill.solid(); arr.fill.fore_color.rgb = NAVY; arr.line.fill.background()
arr.rotation = 90.0

# Fishbone categories (angled lines)
categories = [
    ('PEOPLE', NAVY, ['Weak project\nsponsorship', 'Inadequate\nskills & training', 'Poor\ncommunication']),
    ('PLANNING', TEAL, ['Unrealistic\nschedules', 'Incomplete\nWBS', 'No risk\nallowance']),
    ('PROCUREMENT', GOLD, ['Delayed\nawards', 'Supplier\nfailure', 'Logistics\ndisruption']),
    ('CONTRACTORS', ORANGE, ['Poor\nperformance', 'Financial\ninstability', 'Scope\ndisputes']),
    ('DESIGN', RED, ['Incomplete\nspecs', 'Frequent\nRFIs', 'Late\ndecisions']),
    ('CHANGES', DRED, ['Scope\ncreep', 'Unapproved\nvariations', 'No change\ncontrol']),
]

for i, (cat, color, items) in enumerate(categories):
    angle = -30 if i % 2 == 0 else 30
    y_base = Inches(2.5) if i % 2 == 0 else Inches(4.8)
    x_pos = Inches(1.8) + Inches(i * 1.6)
    
    # Connector line
    R(x_pos, Inches(3.9), Pt(2), Inches(1.8 if i % 2 == 0 else -1.5), color)
    
    # Category box
    R(x_pos - Inches(0.3), y_base, Inches(1.2), Inches(0.4), color)
    T(x_pos - Inches(0.25), y_base + Pt(2), Inches(1.1), Inches(0.35), cat, 9, WHITE, True, PP_ALIGN.CENTER, FLT)
    
    # Cause items
    for j, item in enumerate(items):
        ty = y_base + Inches(0.5) + Inches(j * 0.35)
        T(x_pos - Inches(0.5), ty, Inches(1.5), Inches(0.35), item, 8, RGBColor(60, 60, 60), False, PP_ALIGN.CENTER, FLT)

# Result box
R(Inches(11.3), Inches(3.3), Inches(1.6), Inches(1.0), RED)
T(Inches(11.35), Inches(3.35), Inches(1.5), Inches(0.9), 'PROJECT\nFAILURE', 12, WHITE, True, PP_ALIGN.CENTER, FLT)

# Bottom insight
R(Inches(1.5), Inches(6.0), Inches(10.5), Inches(1.2), RGBColor(245, 246, 248))
T(Inches(1.8), Inches(6.1), Inches(10), Inches(0.4),
  'The Six Failure Drivers  |  Research across 1,200+ failed construction projects (PMI, 2023)',
  12, NAVY, True, PP_ALIGN.LEFT, FLT)
T(Inches(1.8), Inches(6.4), Inches(10), Inches(0.5),
  '72% of projects that fail exhibit 3+ of these drivers simultaneously. The most lethal combination:\nUnrealistic Planning + Poor Change Control + Weak Sponsorship.',
  11, RGBColor(80, 80, 80), False, PP_ALIGN.LEFT, FLT)

# ══════════════════════════════════════════════════════════════
# SLIDE 7  —  COST IMPACT SIMULATION (waterfall chart)
# ══════════════════════════════════════════════════════════════
SL()
R(0, 0, SW, SH, WHITE)
R(0, 0, Inches(0.08), SH, GOLD)

T(Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
  'Cost Impact Simulation  |  How Small Overruns Escalate', 22, NAVY, True, PP_ALIGN.LEFT, FLT)
T(Inches(0.8), Inches(0.75), Inches(11), Inches(0.4),
  'Project Budget: SAR 100 Million  |  Simulating 5%, 10%, 15%, and 20% overrun scenarios', 12, MGRAY, False, PP_ALIGN.LEFT, FLT)

IMG(C('waterfall_cost.png'), Inches(0.3), Inches(1.1), Inches(9), Inches(5.2))

# Right side: cost impact cards
scenarios = [
    ('5%  Overrun', 'SAR 5M', 'Minor impact.\nManageable with\ncontingency.', GOLD),
    ('10%  Overrun', 'SAR 10M', 'Significant.\nRequires scope\nreduction.', ORANGE),
    ('15%  Overrun', 'SAR 15M', 'Critical.\nNeeds capital\ninjection.', RED),
    ('20%  Overrun', 'SAR 20M', 'Catastrophic.\nProject at risk\nof termination.', DRED),
]
for i, (title, amt, impact, color) in enumerate(scenarios):
    y = Inches(1.3) + Inches(i * 1.4)
    R(Inches(9.5), y, Inches(3.5), Inches(1.2), WHITE)
    R(Inches(9.5), y, Inches(0.08), Inches(1.2), color)
    T(Inches(9.8), y + Pt(2), Inches(3.0), Inches(0.3), title, 12, color, True, PP_ALIGN.LEFT, FLT)
    T(Inches(9.8), y + Pt(14), Inches(3.0), Inches(0.3), amt, 18, NAVY, True, PP_ALIGN.LEFT, FLT)
    T(Inches(9.8), y + Pt(32), Inches(3.0), Inches(0.5), impact, 9, MGRAY, False, PP_ALIGN.LEFT, FLT)

# ══════════════════════════════════════════════════════════════
# SLIDE 8  —  SCHEDULE IMPACT (Gantt-style visualization)
# ══════════════════════════════════════════════════════════════
SL()
R(0, 0, SW, SH, DARK)
R(0, 0, Inches(0.08), SH, GOLD)

T(Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
  'Schedule Impact Simulation  |  Delay Scenarios', 24, WHITE, True, PP_ALIGN.LEFT, FLT)
T(Inches(0.8), Inches(0.75), Inches(11), Inches(0.4),
  'How 30, 60, and 90-day delays cascade across project phases', 12, LGRAY, False, PP_ALIGN.LEFT, FLT)

# Gantt-style chart
phases = ['Foundation', 'Structure', 'Envelope', 'MEP', 'Interiors', 'Landscaping', 'Commissioning']
planned_starts = [0, 3, 7, 8, 11, 14, 16]
planned_durs = [3, 4, 3, 4, 4, 3, 2]

for i, (phase, ps, pd) in enumerate(zip(phases, planned_starts, planned_durs)):
    y = Inches(1.5) + Inches(i * 0.65)
    # Phase label
    T(Inches(0.8), y, Inches(1.5), Inches(0.4), phase, 11, WHITE, True, PP_ALIGN.RIGHT)
    # Planned bar
    x_planned = Inches(2.5) + Inches(ps * 0.5)
    w_planned = Inches(pd * 0.5)
    R(x_planned, y, w_planned, Inches(0.25), TEAL, 0.6)
    # Delay bar
    delay = 30 if i < 3 else 60 if i < 5 else 90
    delay_months = delay / 30
    x_delayed = x_planned + w_planned + Inches(0.1)
    w_delayed = Inches(delay_months * 0.5)
    R(x_delayed, y, w_delayed, Inches(0.25), RED, 0.5)
    # Delay label
    T(x_delayed + Inches(0.05), y, w_delayed, Inches(0.25), f'+{delay}d', 8, RED, True, PP_ALIGN.LEFT, FLT)

# Legend
R(Inches(2.5), Inches(6.2), Inches(0.3), Inches(0.2), TEAL, 0.6)
T(Inches(2.9), Inches(6.15), Inches(1.5), Inches(0.25), 'Planned', 9, LGRAY, False, PP_ALIGN.LEFT, FLT)
R(Inches(4.5), Inches(6.2), Inches(0.3), Inches(0.2), RED, 0.5)
T(Inches(4.9), Inches(6.15), Inches(1.5), Inches(0.25), 'Delay', 9, LGRAY, False, PP_ALIGN.LEFT, FLT)

# Cumulative impact
R(Inches(8), Inches(1.5), Inches(4.8), Inches(5.0), NAVY, 0.4)
T(Inches(8.3), Inches(1.6), Inches(4.2), Inches(0.4), 'Cascading Impact', 14, GOLD, True, PP_ALIGN.LEFT)
T(Inches(8.3), Inches(2.1), Inches(4.2), Inches(3.5),
  'A 30-day delay in Foundation\ncompounds to a 90+ day delay\nat Project Completion.\n\nWhy?\n\n•  Each phase depends\n   on the previous\n•  Resources get\n   reallocated\n•  Contractor claims\n   and disputes arise\n•  Liquidated damages\n   accumulate\n•  Financing costs\n   increase',
  10, LGRAY, False, PP_ALIGN.LEFT, FLT)

# ══════════════════════════════════════════════════════════════
# SLIDE 9  —  EXECUTIVE HEAT MAP
# ══════════════════════════════════════════════════════════════
SL()
R(0, 0, SW, SH, WHITE)
R(0, 0, Inches(0.08), SH, GOLD)

T(Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
  'Executive Risk Heat Map  |  Construction Project Exposure', 22, NAVY, True, PP_ALIGN.LEFT, FLT)
T(Inches(0.8), Inches(0.75), Inches(11), Inches(0.4),
  'Impact vs. Probability  |  PMO Risk Assessment Matrix', 12, MGRAY, False, PP_ALIGN.LEFT, FLT)

IMG(C('heatmap_risk.png'), Inches(0.3), Inches(1.2), Inches(9), Inches(5.5))

# Risk register on right
R(Inches(9.5), Inches(1.3), Inches(3.5), Inches(5.5), RGBColor(245, 246, 248))
T(Inches(9.7), Inches(1.4), Inches(3.1), Inches(0.4), 'Top 5 Risks', 14, NAVY, True, PP_ALIGN.LEFT)

risks = [
    ('CRITICAL', 'Contractor financial\ndefault', DRED),
    ('CRITICAL', 'Steel price\nescalation 40%+', DRED),
    ('HIGH', 'Design changes\nby client', RED),
    ('HIGH', 'Labor shortage\n25% gap', RED),
    ('MEDIUM', 'Permit delays\nfrom authority', GOLD),
]
for i, (level, risk, color) in enumerate(risks):
    y = Inches(1.9) + Inches(i * 0.75)
    R(Inches(9.7), y, Inches(3.1), Inches(0.65), WHITE)
    R(Inches(9.7), y, Inches(0.06), Inches(0.65), color)
    T(Inches(9.9), y + Pt(2), Inches(2.7), Inches(0.25), level, 8, color, True, PP_ALIGN.LEFT, FLT)
    T(Inches(9.9), y + Pt(14), Inches(2.7), Inches(0.4), risk, 10, RGBColor(60, 60, 60), False, PP_ALIGN.LEFT, FLT)

# ══════════════════════════════════════════════════════════════
# SLIDE 10  —  CASE STUDY (KPI cards)
# ══════════════════════════════════════════════════════════════
SL()
R(0, 0, SW, SH, WHITE)
R(0, 0, Inches(0.08), SH, GOLD)

T(Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
  'Case Study  |  Mega Airport Terminal Project', 22, NAVY, True, PP_ALIGN.LEFT, FLT)
T(Inches(0.8), Inches(0.75), Inches(11), Inches(0.4),
  'SAR 2.5 Billion  |  48 Months Planned  |  Actual: SAR 3.45B, 68 Months', 12, MGRAY, False, PP_ALIGN.LEFT, FLT)

IMG(C('case_study_kpi.png'), Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.5))

# ══════════════════════════════════════════════════════════════
# SLIDE 11  —  PMO MONITORING DASHBOARD
# ══════════════════════════════════════════════════════════════
SL()
R(0, 0, SW, SH, NAVY)
R(0, 0, Inches(0.08), SH, GOLD)

T(Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
  'Early Detection  |  PMO Performance Monitoring', 24, WHITE, True, PP_ALIGN.LEFT, FLT)
T(Inches(0.8), Inches(0.75), Inches(11), Inches(0.4),
  'How top organizations identify problems before they become crises', 12, LGRAY, False, PP_ALIGN.LEFT, FLT)

IMG(C('pmo_dashboard.png'), Inches(0.3), Inches(1.2), Inches(12.7), Inches(5.8))

# ══════════════════════════════════════════════════════════════
# SLIDE 12  —  THE FUTURE STATE (Before vs After)
# ══════════════════════════════════════════════════════════════
SL()
R(0, 0, SW, SH, WHITE)
R(0, 0, Inches(0.08), SH, GOLD)

T(Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
  'The Future State  |  Transformation Roadmap', 22, NAVY, True, PP_ALIGN.LEFT, FLT)
T(Inches(0.8), Inches(0.75), Inches(11), Inches(0.4),
  'From reactive firefighting to proactive project intelligence', 12, MGRAY, False, PP_ALIGN.LEFT, FLT)

# BEFORE column
R(Inches(0.8), Inches(1.4), Inches(5.5), Inches(5.5), RGBColor(250, 235, 235))
R(Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.06), RED)
T(Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5), 'BEFORE  |  Reactive Organization', 16, RED, True, PP_ALIGN.CENTER, FLT)

before_items = [
    '✗  Monthly reports only',
    '✗  No real-time EVM data',
    '✗  Excel-based tracking',
    '✗  Siloed cost & schedule',
    '✗  Late problem detection',
    '✗  Firefighting culture',
    '✗  No predictive analytics',
    '✗  Ad-hoc decision making',
]
for i, item in enumerate(before_items):
    T(Inches(1.2), Inches(2.2) + Inches(i * 0.5), Inches(4.8), Inches(0.45), item, 13, RGBColor(180, 80, 80), False, PP_ALIGN.LEFT, FLT)

# Arrow
arr = prs.slides[-1].shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.5), Inches(3.8), Inches(0.6), Inches(0.5))
arr.fill.solid(); arr.fill.fore_color.rgb = GOLD; arr.line.fill.background()

# AFTER column
R(Inches(7.2), Inches(1.4), Inches(5.5), Inches(5.5), RGBColor(235, 250, 240))
R(Inches(7.2), Inches(1.4), Inches(5.5), Inches(0.06), GREEN)
T(Inches(7.2), Inches(1.5), Inches(5.5), Inches(0.5), 'AFTER  |  Predictive Organization', 16, GREEN, True, PP_ALIGN.CENTER, FLT)

after_items = [
    '✓  Real-time PMO dashboard',
    '✓  Automated EVM tracking',
    '✓  Integrated P6 + ERP + PMIS',
    '✓  Unified cost/schedule control',
    '✓  Early warning at 20% complete',
    '✓  Predictive risk analytics',
    '✓  AI-powered forecasting',
    '✓  Data-driven decisions',
]
for i, item in enumerate(after_items):
    T(Inches(7.6), Inches(2.2) + Inches(i * 0.5), Inches(4.8), Inches(0.45), item, 13, RGBColor(40, 130, 70), False, PP_ALIGN.LEFT, FLT)

# ══════════════════════════════════════════════════════════════
# SLIDE 13  —  KEY RECOMMENDATIONS (visual strategy map)
# ══════════════════════════════════════════════════════════════
SL()
R(0, 0, SW, SH, WHITE)
R(0, 0, Inches(0.08), SH, GOLD)

T(Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
  'Strategic Recommendations  |  The Path Forward', 22, NAVY, True, PP_ALIGN.LEFT, FLT)
T(Inches(0.8), Inches(0.75), Inches(11), Inches(0.4),
  'Five actions that will transform your project controls capability', 12, MGRAY, False, PP_ALIGN.LEFT, FLT)

# Strategy map with connected boxes
recs = [
    ('01', 'Implement\nEVM', 'Deploy Earned Value\nManagement across\nall projects > SAR 50M', TEAL),
    ('02', 'Integrate\nSystems', 'Connect P6, ERP,\nand PMIS into a\nsingle data platform', NAVY),
    ('03', 'Build PMO\nCapability', 'Establish a dedicated\nProject Controls PMO\nwith EVM expertise', GOLD),
    ('04', 'Automate\nReporting', 'Real-time dashboards\nwith automated\nearly warning alerts', TEAL),
    ('05', 'Forecast\nProactively', 'Use predictive analytics\nfor cost and schedule\nforecasting', NAVY),
]

for i, (num, title, desc, color) in enumerate(recs):
    x = Inches(0.5) + Inches(i * 2.55)
    
    # Connector line (except first)
    if i > 0:
        R(x - Inches(0.1), Inches(3.5), Inches(0.2), Pt(3), GOLD)
    
    # Circle with number
    circ = prs.slides[-1].shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.7), Inches(1.5), Inches(0.6), Inches(0.6))
    circ.fill.solid(); circ.fill.fore_color.rgb = color; circ.line.fill.background()
    T(x + Inches(0.75), Inches(1.55), Inches(0.5), Inches(0.5), num, 18, WHITE, True, PP_ALIGN.CENTER, FLT)
    
    # Title
    T(x + Inches(0.1), Inches(2.2), Inches(2.0), Inches(0.6), title, 16, color, True, PP_ALIGN.CENTER, FLT)
    
    # Description box
    R(x, Inches(3.0), Inches(2.3), Inches(1.5), color, 0.08)
    T(x + Inches(0.15), Inches(3.1), Inches(2.0), Inches(1.2), desc, 11, RGBColor(60, 60, 60), False, PP_ALIGN.CENTER, FLT)

# Bottom implementation timeline
R(Inches(0.5), Inches(5.0), Inches(12.3), Inches(2), RGBColor(245, 246, 248))
T(Inches(0.8), Inches(5.1), Inches(11.7), Inches(0.3), 'Implementation Timeline', 13, NAVY, True, PP_ALIGN.CENTER)

phases = [
    ('Month 1-3', 'Phase 1:\nFoundation', 'EVM setup,\nWBS, baseline', TEAL),
    ('Month 4-6', 'Phase 2:\nIntegration', 'P6 + ERP\nconnection', NAVY),
    ('Month 7-9', 'Phase 3:\nPMO Launch', 'Dashboard,\ntraining, ops', GOLD),
    ('Month 10-12', 'Phase 4:\nOptimization', 'Advanced analytics,\nAI forecasting', TEAL),
]
for i, (time, phase, desc, color) in enumerate(phases):
    x = Inches(1.2) + Inches(i * 3.0)
    R(x, Inches(5.6), Inches(2.6), Inches(1.3), WHITE)
    R(x, Inches(5.6), Inches(2.6), Inches(0.06), color)
    T(x + Inches(0.1), Inches(5.65), Inches(2.4), Inches(0.25), time, 9, color, True, PP_ALIGN.CENTER, FLT)
    T(x + Inches(0.1), Inches(5.85), Inches(2.4), Inches(0.4), phase, 12, NAVY, True, PP_ALIGN.CENTER, FLT)
    T(x + Inches(0.1), Inches(6.2), Inches(2.4), Inches(0.4), desc, 9, MGRAY, False, PP_ALIGN.CENTER, FLT)
    if i < len(phases) - 1:
        R(x + Inches(2.6), Inches(6.2), Inches(0.4), Pt(2), GOLD)

# ══════════════════════════════════════════════════════════════
# SLIDE 14  —  FINAL EXECUTIVE MESSAGE
# ══════════════════════════════════════════════════════════════
SL()
IMG(K('0000_Overlook_View_1920x1080.jpg'))
R(0, 0, SW, SH, DARK, 0.55)
R(0, 0, Inches(0.12), SH, GOLD)

# Center content
R(Inches(4.5), Inches(2.2), Inches(4.3), Pt(4), GOLD)

TX(Inches(1.5), Inches(2.5), Inches(10.3), Inches(3), [
    ('CONSTRUCTION PROJECT FAILURE', 14, LGRAY, False, FLT),
    ('IS NOT INEVITABLE', 40, WHITE, True, FLT),
    ('', 12, WHITE, False, FLT),
    ('The data is clear. The tools exist. The methodology is proven.', 16, GOLD, False, FLT),
    ('', 8, WHITE, False, FLT),
    ('What separates successful organizations from failing ones is not budget or size —', 13, LGRAY, False, FLT),
    ('it is the discipline to detect problems early and act decisively.', 13, LGRAY, False, FLT),
], PP_ALIGN.CENTER)

# Bottom bar
R(Inches(2), Inches(6.0), Inches(9.3), Inches(0.6), NAVY, 0.7)
T(Inches(2.2), Inches(6.02), Inches(8.9), Inches(0.5),
  'PMO Advisory  |  Project Controls Excellence  |  Data-Driven Project Management',
  12, GOLD, False, PP_ALIGN.CENTER, FLT)

# ══════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
print(f"Slides: {len(prs.slides)}")
