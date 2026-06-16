from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
from pptx.oxml.ns import qn
import os

IMG_DIR = r"E:\N8N\scraper\scraper2\csi-ultimate\ksp_images"
CHART_DIR = r"E:\N8N\scraper\scraper2\csi-ultimate\evm_charts"
OUTPUT = r"E:\N8N\scraper\scraper2\csi-ultimate\EVM_Construction_Projects.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SW = prs.slide_width
SH = prs.slide_height

# Colors
NAVY = RGBColor(27, 42, 74)
DARK_NAVY = RGBColor(18, 30, 55)
GOLD = RGBColor(200, 150, 46)
TEAL = RGBColor(26, 138, 158)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(200, 200, 210)
MED_GRAY = RGBColor(150, 150, 165)
DARK_BG = RGBColor(15, 25, 50)
RED = RGBColor(192, 57, 43)
GREEN = RGBColor(39, 174, 96)
ORANGE = RGBColor(230, 126, 34)
SOFT_WHITE = RGBColor(240, 240, 245)

FONT = 'Calibri'
FONT_LIGHT = 'Calibri Light'

def set_alpha(shape, alpha_val):
    spPr = shape._element.spPr
    sf = spPr.find(qn('a:solidFill'))
    if sf is not None:
        srgb = sf.find(qn('a:srgbClr'))
        if srgb is not None:
            srgb.set('alpha', str(int(alpha_val * 1000)))

def add_shape(slide, left, top, width, height, fill_color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if alpha is not None:
        set_alpha(shape, alpha)
    return shape

def add_full_bg(slide, image_path):
    slide.shapes.add_picture(image_path, 0, 0, SW, SH)

def add_txt(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font=FONT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return tb

def add_rich(slide, left, top, width, height, lines, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (txt, sz, clr, bld, fn) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(sz)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.font.name = fn
        p.alignment = align
        p.space_after = Pt(4)
    return tb

def add_gold_accent(slide, left, top, width=Inches(0.8), height=Pt(3)):
    add_shape(slide, left, top, width, height, GOLD)

def slide_section(slide, title, subtitle=None):
    add_shape(slide, 0, 0, SW, SH, NAVY, alpha=0.0)
    add_gold_accent(slide, Inches(0.8), Inches(0.6), Inches(0.6), Pt(4))
    add_txt(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.6), title, 32, WHITE, True, PP_ALIGN.LEFT, FONT_LIGHT)
    if subtitle:
        add_txt(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.4), subtitle, 14, GOLD, False, PP_ALIGN.LEFT, FONT_LIGHT)

def slide_title_bar(slide, title, title_ar):
    add_shape(slide, 0, 0, SW, Inches(1.1), NAVY)
    add_shape(slide, 0, Inches(1.1), SW, Pt(3), GOLD)
    add_txt(slide, Inches(0.8), Inches(0.15), Inches(11), Inches(0.8), title_ar, 26, WHITE, True, PP_ALIGN.RIGHT, FONT)
    add_txt(slide, Inches(0.8), Inches(0.15), Inches(11), Inches(0.8), title, 14, GOLD, False, PP_ALIGN.LEFT, FONT_LIGHT)

def make_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])

def img(name):
    return os.path.join(IMG_DIR, name)

def chart(name):
    return os.path.join(CHART_DIR, name)

# =====================================================================
# SLIDE 1: COVER
# =====================================================================
s = make_slide()
add_full_bg(s, img('1509_KSP_Park_Overview_001_1920x1080.jpg'))
add_shape(s, 0, 0, SW, SH, DARK_BG, 0.6)
add_shape(s, 0, 0, Inches(0.12), SH, GOLD)
add_gold_accent(s, Inches(1.5), Inches(2.0), Inches(1.5), Pt(5))
add_rich(s, Inches(1.5), Inches(2.3), Inches(10), Inches(4), [
    ('إدارة القيمة المكتسبة', 44, WHITE, True, FONT_LIGHT),
    ('Earned Value Management', 20, GOLD, False, FONT_LIGHT),
    ('في المشاريع الإنشائية', 36, WHITE, True, FONT_LIGHT),
    ('', 12, WHITE, False, FONT_LIGHT),
    ('دليل شامل لتطبيق منهجية EVM في إدارة المشاريع', 16, LIGHT_GRAY, False, FONT_LIGHT),
    ('Engineering  |  Construction  |  Infrastructure  |  Mega Projects', 13, MED_GRAY, False, FONT_LIGHT),
], PP_ALIGN.LEFT)
add_rich(s, Inches(1.5), Inches(6.0), Inches(10), Inches(1), [
    ('PMO  |  Project Controls  |  Cost Management  |  PMI-PMP', 12, MED_GRAY, False, FONT_LIGHT),
    ('برنامج تدريبي احترافي  |  النسخة 1.0', 11, MED_GRAY, False, FONT_LIGHT),
], PP_ALIGN.LEFT)

# =====================================================================
# SLIDE 2: AGENDA
# =====================================================================
s = make_slide()
add_shape(s, 0, 0, SW, SH, NAVY)
add_shape(s, 0, 0, Inches(0.12), SH, GOLD)
slide_section(s, 'Agenda  |  جدول الأعمال')
items_agenda = [
    ('1', 'مقدمة عن إدارة القيمة المكتسبة', 'Introduction to EVM'),
    ('2', 'تاريخ EVM ومعايير PMI', 'History & PMI Standards'),
    ('3', 'مكونات EVM الأساسية', 'Core EVM Components'),
    ('4', 'المعادلات والمؤشرات الرئيسية', 'Key Formulas & Indicators'),
    ('5', 'منحنيات S وتحليل الأداء', 'S-Curves & Performance Analysis'),
    ('6', 'تطبيقات EVM في المشاريع الإنشائية', 'Construction Applications'),
    ('7', 'دراسة حالة عملية', 'Real-World Case Study'),
    ('8', 'لوحات القيادة التنفيذية', 'Executive Dashboards'),
    ('9', 'التكامل مع الأنظمة', 'EVM Integration'),
    ('10', 'أفضل الممارسات والتوصيات', 'Best Practices'),
]
for i, (num, ar_title, en_title) in enumerate(items_agenda):
    y = Inches(1.4) + Inches(i * 0.55)
    add_shape(s, Inches(0.8), y, Inches(0.5), Inches(0.4), GOLD, 0.3)
    add_txt(s, Inches(0.85), y + Pt(2), Inches(0.4), Inches(0.35), num, 16, GOLD, True, PP_ALIGN.CENTER)
    add_txt(s, Inches(1.5), y + Pt(1), Inches(6), Inches(0.35), ar_title, 15, WHITE, False, PP_ALIGN.RIGHT)
    add_txt(s, Inches(7.5), y + Pt(3), Inches(5), Inches(0.3), en_title, 11, MED_GRAY, False, PP_ALIGN.LEFT, FONT_LIGHT)

# =====================================================================
# SLIDE 3: EXECUTIVE SUMMARY
# =====================================================================
s = make_slide()
add_full_bg(s, img('KSP_Overview.jpg'))
add_shape(s, 0, 0, SW, SH, DARK_BG, 0.65)
slide_title_bar(s, 'Executive Summary', 'ملخص تنفيذي')
add_rich(s, Inches(0.8), Inches(1.4), Inches(5.5), Inches(5.5), [
    ('ما هي إدارة القيمة المكتسبة؟', 22, GOLD, True, FONT),
    ('', 8, WHITE, False, FONT),
    ('منهجية إدارة مشاريع متكاملة تجمع بين نطاق العمل،', 15, WHITE, False, FONT),
    ('الجدول الزمني، والتكاليف في إطار موحد لقياس', 15, WHITE, False, FONT),
    ('أداء المشروع بدقة وموضوعية.', 15, WHITE, False, FONT),
    ('', 12, WHITE, False, FONT),
    ('تتيح EVM للمدراء الإجابة على:', 16, GOLD, True, FONT),
    ('', 6, WHITE, False, FONT),
    ('• هل نحن متقدمون أم متأخرون عن الجدول؟', 14, WHITE, False, FONT),
    ('• هل نحن في حدود الميزانية أم تجاوزناها؟', 14, WHITE, False, FONT),
    ('• كم سيكلف المشروع عند الإكمال؟', 14, WHITE, False, FONT),
    ('• متى سننهي المشروع؟', 14, WHITE, False, FONT),
], PP_ALIGN.RIGHT)
# Key metrics box
add_shape(s, Inches(7), Inches(1.5), Inches(5.5), Inches(5.2), NAVY, 0.4)
add_txt(s, Inches(7.3), Inches(1.7), Inches(5), Inches(0.4), 'EVM Core Benefits', 16, GOLD, True, PP_ALIGN.LEFT, FONT_LIGHT)
benefits = [
    '✓  Objective performance measurement',
    '✓  Early warning indicators',
    '✓  Data-driven decision making',
    '✓  Integrated cost & schedule control',
    '✓  Accurate forecasting',
    '✓  PMI & ISO standard methodology',
    '✓  Applicable to all project sizes',
    '✓  International best practice',
]
for i, b in enumerate(benefits):
    add_txt(s, Inches(7.3), Inches(2.2) + Inches(i * 0.45), Inches(5), Inches(0.4), b, 13, WHITE, False, PP_ALIGN.LEFT, FONT_LIGHT)

# =====================================================================
# SLIDE 4: WHAT IS EVM
# =====================================================================
s = make_slide()
add_shape(s, 0, 0, SW, SH, NAVY)
add_shape(s, 0, 0, Inches(0.12), SH, GOLD)
slide_section(s, 'Introduction to EVM', 'ما هي إدارة القيمة المكتسبة؟')
add_rich(s, Inches(0.8), Inches(1.6), Inches(6), Inches(5), [
    ('مفهوم القيمة المكتسبة', 20, GOLD, True, FONT),
    ('', 8, WHITE, False, FONT),
    ('Earned Value Management (EVM) هي منهجية ', 15, WHITE, False, FONT),
    ('إدارة مشاريع تستخدم ثلاثة أبعاد أساسية:', 15, WHITE, False, FONT),
    ('', 6, WHITE, False, FONT),
    ('1. القيمة المخططة  - Planned Value (PV)', 15, TEAL, True, FONT),
    ('الميزانية المخصصة للعمل المجدول', 13, LIGHT_GRAY, False, FONT),
    ('', 6, WHITE, False, FONT),
    ('2. القيمة المكتسبة  - Earned Value (EV)', 15, GOLD, True, FONT),
    ('قيمة العمل المنجز فعلياً', 13, LIGHT_GRAY, False, FONT),
    ('', 6, WHITE, False, FONT),
    ('3. التكلفة الفعلية  - Actual Cost (AC)', 15, RED, True, FONT),
    ('التكلفة الحقيقية للعمل المنجز', 13, LIGHT_GRAY, False, FONT),
], PP_ALIGN.RIGHT)
# Three pillars box
add_shape(s, Inches(7.2), Inches(1.6), Inches(5.3), Inches(5), DARK_NAVY, 0.5)
add_txt(s, Inches(7.5), Inches(1.8), Inches(4.8), Inches(0.4), 'Why EVM Matters', 18, GOLD, True, PP_ALIGN.LEFT, FONT_LIGHT)
pillars = [
    ('Visibility', 'Real-time project health monitoring'),
    ('Accountability', 'Clear ownership of cost & schedule'),
    ('Predictability', 'Accurate forecasts and early warnings'),
    ('Control', 'Proactive rather than reactive management'),
    ('Governance', 'PMI, ISO 21500, and ANSI/EIA-748 compliant'),
]
for i, (t, d) in enumerate(pillars):
    y = Inches(2.4) + Inches(i * 0.85)
    add_shape(s, Inches(7.5), y, Inches(4.8), Inches(0.75), TEAL, 0.15)
    add_txt(s, Inches(7.7), y + Pt(2), Inches(4.5), Inches(0.35), t, 15, TEAL, True, PP_ALIGN.LEFT)
    add_txt(s, Inches(7.7), y + Pt(22), Inches(4.5), Inches(0.35), d, 12, LIGHT_GRAY, False, PP_ALIGN.LEFT, FONT_LIGHT)

# =====================================================================
# SLIDE 5: HISTORY
# =====================================================================
s = make_slide()
add_shape(s, 0, 0, SW, SH, WHITE)
add_shape(s, 0, 0, Inches(0.12), SH, GOLD)
slide_section(s, 'History & PMI Standards', 'التاريخ والمعايير')
add_shape(s, 0, Inches(1.2), SW, Pt(2), RGBColor(230, 230, 230))

timeline = [
    ('1960s', 'ظهور EVM في وزارة الدفاع الأمريكية (DoD)'),
    ('1967', 'تطوير نظام C/SCSC للرقابة على التكاليف'),
    ('1996', 'اعتماد ANSI/EIA-748 كمعيار أمريكي'),
    ('2000', 'دمج EVM في معيار PMI PMBOK Guide'),
    ('2005', 'إصدار معيار PMI Practice Standard for EVM'),
    ('2011', 'إصدار PMI Standard for EVM (الطبعة الثانية)'),
    ('2019', 'دمج EVM في كل دورات PMP و PfMP'),
    ('2024', 'تطبيق EVM في المشاريع العملاقة عالمياً'),
]
for i, (year, desc) in enumerate(timeline):
    y = Inches(1.5) + Inches(i * 0.65)
    add_shape(s, Inches(0.8), y, Inches(1.3), Inches(0.45), NAVY if i % 2 == 0 else TEAL)
    add_txt(s, Inches(0.85), y + Pt(4), Inches(1.2), Inches(0.35), year, 14, WHITE, True, PP_ALIGN.CENTER)
    add_txt(s, Inches(2.3), y + Pt(4), Inches(5), Inches(0.35), desc, 14, NAVY if i % 2 == 0 else RGBColor(60, 60, 60), False, PP_ALIGN.RIGHT)
    if i < len(timeline) - 1:
        add_shape(s, Inches(1.65), y + Inches(0.5), Inches(0.02), Inches(0.15), GOLD)

# PMI Standards box
add_shape(s, Inches(7.2), Inches(1.5), Inches(5.3), Inches(5.2), NAVY)
add_txt(s, Inches(7.5), Inches(1.7), Inches(4.8), Inches(0.4), 'PMI Standards for EVM', 16, GOLD, True, PP_ALIGN.LEFT)
pmi_items = [
    'PMBOK Guide (7th Edition)',
    'Practice Standard for EVM',
    'Standard for Project Planning',
    'Standard for Risk Management',
    'PMI-PMP Examination Content',
    'PMI-SP Scheduling Standard',
    'PMI-RMP Risk Standard',
]
for i, item in enumerate(pmi_items):
    add_shape(s, Inches(7.5), Inches(2.3) + Inches(i * 0.55), Inches(4.8), Inches(0.45), TEAL, 0.15)
    add_txt(s, Inches(7.7), Inches(2.35) + Inches(i * 0.55), Inches(4.4), Inches(0.35), f'  {item}', 13, WHITE, False, PP_ALIGN.LEFT, FONT_LIGHT)

# =====================================================================
# SLIDE 6: EVM FRAMEWORK
# =====================================================================
s = make_slide()
add_shape(s, 0, 0, SW, SH, WHITE)
add_shape(s, 0, 0, Inches(0.12), SH, GOLD)
slide_section(s, 'EVM Framework', 'إطار عمل القيمة المكتسبة')

# Framework boxes with connectors
framework = [
    ('Work\nBreakdown\nStructure', NAVY, 'تقسيم العمل إلى حزم\nصغيرة قابلة للقياس'),
    ('Planned\nValue (PV)', TEAL, 'تحديد الميزانية\nلكل حزمة عمل'),
    ('Performance\nMeasurement\nBaseline', GOLD, 'وضع خط الأساس\nلقياس الأداء'),
    ('Earned\nValue (EV)', TEAL, 'قياس العمل\nالمنجز فعلياً'),
    ('Performance\nReports', NAVY, 'تقارير الأداء\nوالمؤشرات'),
]
for i, (title, color, sub) in enumerate(framework):
    x = Inches(0.8) + Inches(i * 2.5)
    y = Inches(2.5)
    add_shape(s, x, y, Inches(2.0), Inches(2.0), color, 0.9)
    add_txt(s, x + Inches(0.15), y + Inches(0.3), Inches(1.7), Inches(1.0), title, 16, WHITE, True, PP_ALIGN.CENTER)
    add_txt(s, x + Inches(0.15), y + Inches(1.2), Inches(1.7), Inches(0.6), sub, 11, RGBColor(230, 230, 230), False, PP_ALIGN.CENTER, FONT_LIGHT)
    if i < len(framework) - 1:
        add_shape(s, x + Inches(2.0), y + Inches(0.85), Inches(0.5), Pt(3), GOLD)

# Bottom description
add_shape(s, Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.5), DARK_NAVY)
add_txt(s, Inches(1.2), Inches(5.4), Inches(10.5), Inches(1.2),
    'يقوم إطار EVM على دمج نطاق العمل (WBS) مع الجدول الزمني والتكاليف\nفي خط أساس واحد لقياس الأداء، مما يتيح رؤية متكاملة لحالة المشروع',
    13, LIGHT_GRAY, False, PP_ALIGN.CENTER, FONT_LIGHT)

# =====================================================================
# SLIDE 7: BENEFITS
# =====================================================================
s = make_slide()
add_full_bg(s, img('1509_KSP_Full_Park_001_1920x1080.jpg'))
add_shape(s, 0, 0, SW, SH, DARK_BG, 0.6)
slide_title_bar(s, 'Benefits of EVM', 'فوائد إدارة القيمة المكتسبة')

benefits_data = [
    ('رقابة متكاملة', 'Integration', 'دمج نطاق العمل، الجدول الزمني، والتكاليف', TEAL),
    ('إنذار مبكر', 'Early Warning', 'الكشف عن الانحرافات قبل تفاقمها', GOLD),
    ('تنبؤ دقيق', 'Accurate Forecast', 'تقدير دقيق لتكلفة ومدة الإنجاز', NAVY),
    ('شفافية كاملة', 'Full Transparency', 'مقاييس موضوعية وواضحة لجميع أصحاب المصلحة', TEAL),
    ('اتخاذ القرارات', 'Decision Support', 'قرارات مبنية على بيانات وليس حدساً', GOLD),
]
for i, (ar, en, desc, color) in enumerate(benefits_data):
    x = Inches(0.8) + Inches(i * 2.5)
    add_shape(s, x, Inches(1.5), Inches(2.2), Inches(3.5), color, 0.15)
    add_shape(s, x, Inches(1.5), Inches(2.2), Inches(0.8), color, 0.8)
    add_txt(s, x + Inches(0.1), Inches(1.55), Inches(2.0), Inches(0.35), ar, 16, WHITE, True, PP_ALIGN.CENTER)
    add_txt(s, x + Inches(0.1), Inches(1.85), Inches(2.0), Inches(0.3), en, 10, RGBColor(220, 220, 220), False, PP_ALIGN.CENTER, FONT_LIGHT)
    add_txt(s, x + Inches(0.15), Inches(2.5), Inches(1.9), Inches(2.0), desc, 12, WHITE, False, PP_ALIGN.CENTER, FONT_LIGHT)

add_txt(s, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1),
    'تطبق منهجية EVM في المشاريع الإنشائية والبنية التحتية والنفط والغاز والطاقة والمشاريع الحكومية',
    14, LIGHT_GRAY, False, PP_ALIGN.CENTER, FONT_LIGHT)

# =====================================================================
# SLIDE 8: PV - Planned Value
# =====================================================================
s = make_slide()
add_shape(s, 0, 0, SW, SH, WHITE)
add_shape(s, 0, 0, Inches(0.12), SH, GOLD)
slide_section(s, 'Planned Value (PV)', 'القيمة المخططة')

add_shape(s, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.2), NAVY)
add_rich(s, Inches(1.2), Inches(1.7), Inches(5), Inches(4.5), [
    ('القيمة المخططة  PV', 22, GOLD, True, FONT),
    ('', 10, WHITE, False, FONT),
    ('التعريف:', 15, TEAL, True, FONT),
    ('الميزانية المعتمدة للعمل المقرر إنجازه', 14, WHITE, False, FONT),
    ('ضمن جدول زمني محدد.', 14, WHITE, False, FONT),
    ('', 8, WHITE, False, FONT),
    ('المعادلة:', 15, GOLD, True, FONT),
    ('PV = BAC × (% من العمل المجدول)', 18, WHITE, True, FONT),
    ('', 8, WHITE, False, FONT),
    ('مثال عملي:', 15, TEAL, True, FONT),
    ('مشروع بناء بمساحة 10,000 متر مربع', 14, WHITE, False, FONT),
    ('الميزانية: 100 مليون ريال', 14, WHITE, False, FONT),
    ('المدة: 12 شهراً', 14, WHITE, False, FONT),
    ('PV بعد 8 أشهر = 78 مليون ريال', 16, GOLD, True, FONT),
], PP_ALIGN.RIGHT)

# Formula visualization box
add_shape(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.2), RGBColor(245, 245, 250))
add_txt(s, Inches(7.1), Inches(1.7), Inches(5.2), Inches(0.4), 'PV at a Glance', 16, NAVY, True, PP_ALIGN.LEFT)
pv_detail = [
    'Also known as: Budgeted Cost of Work Scheduled (BCWS)',
    '',
    'Key Characteristics:',
    '  • Authorized budget for scheduled work',
    '  • Time-phased expenditure plan',
    '  • Forms the Performance Measurement Baseline',
    '  • Measured in currency (SAR, USD, etc.)',
    '',
    'Construction Application:',
    '  • Monthly payment applications',
    '  • Contractor progress billing',
    '  • Resource-loaded schedule',
    '  • Cash flow forecasting',
]
for i, line in enumerate(pv_detail):
    c = NAVY if line and line[0].isupper() else RGBColor(80, 80, 80)
    b = bool(line and line[0].isupper())
    add_txt(s, Inches(7.1), Inches(2.2) + Inches(i * 0.3), Inches(5.2), Inches(0.28), line, 11, c, b, PP_ALIGN.LEFT, FONT_LIGHT)

# =====================================================================
# SLIDE 9: EV - Earned Value
# =====================================================================
s = make_slide()
add_shape(s, 0, 0, SW, SH, WHITE)
add_shape(s, 0, 0, Inches(0.12), SH, GOLD)
slide_section(s, 'Earned Value (EV)', 'القيمة المكتسبة')

add_shape(s, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.2), NAVY)
add_rich(s, Inches(1.2), Inches(1.7), Inches(5), Inches(4.5), [
    ('القيمة المكتسبة  EV', 22, GOLD, True, FONT),
    ('', 10, WHITE, False, FONT),
    ('التعريف:', 15, TEAL, True, FONT),
    ('قيمة العمل المنجز فعلياً مقاسة', 14, WHITE, False, FONT),
    ('بالميزانية المعتمدة.', 14, WHITE, False, FONT),
    ('', 8, WHITE, False, FONT),
    ('المعادلة:', 15, GOLD, True, FONT),
    ('EV = BAC × (% من العمل المنجز)', 18, WHITE, True, FONT),
    ('', 8, WHITE, False, FONT),
    ('مثال عملي:', 15, TEAL, True, FONT),
    ('تم إنجاز 75% من الأعمال الإنشائية', 14, WHITE, False, FONT),
    ('EV = 100 × 75% = 75 مليون ريال', 18, GOLD, True, FONT),
], PP_ALIGN.RIGHT)

add_shape(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.2), RGBColor(245, 245, 250))
add_txt(s, Inches(7.1), Inches(1.7), Inches(5.2), Inches(0.4), 'EV at a Glance', 16, NAVY, True, PP_ALIGN.LEFT)
ev_detail = [
    'Also known as: Budgeted Cost of Work Performed (BCWP)',
    '',
    'Key Characteristics:',
    '  • Represents physical progress in monetary terms',
    '  • Based on measurable completion criteria',
    '  • Independent of actual costs incurred',
    '  • The "earned" value of work done',
    '',
    'Construction Application:',
    '  • Physical % complete × Budget',
    '  • Milestone completion measurement',
    '  • Units completed method',
    '  • Weighted milestone method',
]
for i, line in enumerate(ev_detail):
    c = NAVY if line and line[0].isupper() else RGBColor(80, 80, 80)
    b = bool(line and line[0].isupper())
    add_txt(s, Inches(7.1), Inches(2.2) + Inches(i * 0.3), Inches(5.2), Inches(0.28), line, 11, c, b, PP_ALIGN.LEFT, FONT_LIGHT)

# =====================================================================
# SLIDE 10: AC - Actual Cost
# =====================================================================
s = make_slide()
add_shape(s, 0, 0, SW, SH, WHITE)
add_shape(s, 0, 0, Inches(0.12), SH, GOLD)
slide_section(s, 'Actual Cost (AC)', 'التكلفة الفعلية')

add_shape(s, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.2), NAVY)
add_rich(s, Inches(1.2), Inches(1.7), Inches(5), Inches(4.5), [
    ('التكلفة الفعلية  AC', 22, GOLD, True, FONT),
    ('', 10, WHITE, False, FONT),
    ('التعريف:', 15, TEAL, True, FONT),
    ('إجمالي التكاليف المحققة فعلياً', 14, WHITE, False, FONT),
    ('للعمل المنجز حتى تاريخ التقرير.', 14, WHITE, False, FONT),
    ('', 8, WHITE, False, FONT),
    ('المعادلة:', 15, GOLD, True, FONT),
    ('AC = مجموع التكاليف الفعلية', 18, WHITE, True, FONT),
    ('(مباشرة + غير مباشرة)', 14, LIGHT_GRAY, False, FONT),
    ('', 8, WHITE, False, FONT),
    ('مثال عملي:', 15, TEAL, True, FONT),
    ('بلغت التكاليف الفعلية بعد 8 أشهر', 14, WHITE, False, FONT),
    ('AC = 82 مليون ريال', 18, GOLD, True, FONT),
], PP_ALIGN.RIGHT)

add_shape(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.2), RGBColor(245, 245, 250))
add_txt(s, Inches(7.1), Inches(1.7), Inches(5.2), Inches(0.4), 'AC at a Glance', 16, NAVY, True, PP_ALIGN.LEFT)
ac_detail = [
    'Also known as: Actual Cost of Work Performed (ACWP)',
    '',
    'Key Characteristics:',
    '  • Total costs incurred for completed work',
    '  • Includes labor, materials, equipment, subcontractors',
    '  • Captured from accounting/ERP systems',
    '  • Time-phased for accurate tracking',
    '',
    'Cost Tracking Methods:',
    '  • Direct costs (labor, materials, equipment)',
    '  • Indirect costs (overhead, G&A)',
    '  • Commitments and accruals',
    '  • Progress payment applications',
]
for i, line in enumerate(ac_detail):
    c = NAVY if line and line[0].isupper() else RGBColor(80, 80, 80)
    b = bool(line and line[0].isupper())
    add_txt(s, Inches(7.1), Inches(2.2) + Inches(i * 0.3), Inches(5.2), Inches(0.28), line, 11, c, b, PP_ALIGN.LEFT, FONT_LIGHT)

# =====================================================================
# SLIDE 11: BAC
# =====================================================================
s = make_slide()
add_shape(s, 0, 0, SW, SH, WHITE)
add_shape(s, 0, 0, Inches(0.12), SH, GOLD)
slide_section(s, 'Budget at Completion (BAC)', 'الميزانية عند الإكمال')

add_shape(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.2), NAVY, 0.05)
add_rich(s, Inches(1.2), Inches(1.7), Inches(5), Inches(4.5), [
    ('BAC  الميزانية عند الإكمال', 22, GOLD, True, FONT),
    ('', 10, WHITE, False, FONT),
    ('التعريف:', 15, TEAL, True, FONT),
    ('إجمالي الميزانية المعتمدة للمشروع.', 14, WHITE, False, FONT),
    ('', 8, WHITE, False, FONT),
    ('BAC هو الأساس الذي تقاس عليه:', 15, GOLD, True, FONT),
    ('• القيمة المخططة (PV)', 14, WHITE, False, FONT),
    ('• القيمة المكتسبة (EV)', 14, WHITE, False, FONT),
    ('• جميع انحرافات التكلفة والجدول', 14, WHITE, False, FONT),
    ('', 8, WHITE, False, FONT),
    ('BAC = إجمالي ميزانية المشروع', 18, GOLD, True, FONT),
])
add_rich(s, Inches(7), Inches(1.7), Inches(5), Inches(4.5), [
    ('BAC in Practice', 18, TEAL, True, FONT_LIGHT),
    ('', 10, WHITE, False, FONT),
    ('BAC represents the total budget for all authorized work.', 12, LIGHT_GRAY, False, FONT_LIGHT),
    ('', 6, WHITE, False, FONT),
    ('In construction projects, BAC includes:', 13, WHITE, True, FONT_LIGHT),
    ('  • Direct costs (labor, materials, equipment)', 12, LIGHT_GRAY, False, FONT_LIGHT),
    ('  • Indirect costs (site overhead, HO overhead)', 12, LIGHT_GRAY, False, FONT_LIGHT),
    ('  • Profit margin (for contractors)', 12, LIGHT_GRAY, False, FONT_LIGHT),
    ('  • Contingency allowance', 12, LIGHT_GRAY, False, FONT_LIGHT),
    ('  • Escalation / inflation provisions', 12, LIGHT_GRAY, False, FONT_LIGHT),
    ('', 6, WHITE, False, FONT),
    ('BAC is established during the planning phase', 12, LIGHT_GRAY, False, FONT_LIGHT),
    ('and approved as part of the project baseline.', 12, LIGHT_GRAY, False, FONT_LIGHT),
    ('', 6, WHITE, False, FONT),
    ('Changes to BAC require formal change control', 12, LIGHT_GRAY, False, FONT_LIGHT),
    ('through the approved change management process.', 12, LIGHT_GRAY, False, FONT_LIGHT),
])

# =====================================================================
# SLIDE 12: THREE PILLARS
# =====================================================================
s = make_slide()
add_full_bg(s, img('1509_KSP_006_1920x1080.jpg'))
add_shape(s, 0, 0, SW, SH, DARK_BG, 0.6)
slide_title_bar(s, 'The Three Pillars of EVM', 'الأركان الثلاثة لإدارة القيمة المكتسبة')

pillars_data = [
    ('PV', 'Planned\nValue', 'القيمة\nالمخططة', 'أين كنا\nنتوقع أن نكون؟', NAVY),
    ('EV', 'Earned\nValue', 'القيمة\nالمكتسبة', 'أين نحن\nفعلياً؟', TEAL),
    ('AC', 'Actual\nCost', 'التكلفة\nالفعلية', 'كم أنفقنا\nحتى الآن؟', GOLD),
]
for i, (abbr, en, ar, q, color) in enumerate(pillars_data):
    x = Inches(0.8) + Inches(i * 4.1)
    add_shape(s, x, Inches(1.6), Inches(3.8), Inches(5), color, 0.15)
    add_shape(s, x, Inches(1.6), Inches(3.8), Inches(1.2), color, 0.8)
    add_txt(s, x, Inches(1.65), Inches(3.8), Inches(0.5), abbr, 32, WHITE, True, PP_ALIGN.CENTER)
    add_txt(s, x, Inches(2.1), Inches(3.8), Inches(0.3), en, 12, RGBColor(220, 220, 220), False, PP_ALIGN.CENTER, FONT_LIGHT)
    add_txt(s, x + Inches(0.3), Inches(3.0), Inches(3.2), Inches(1), ar, 24, WHITE, True, PP_ALIGN.CENTER)
    add_txt(s, x + Inches(0.3), Inches(4.2), Inches(3.2), Inches(0.8), q, 14, RGBColor(200, 200, 200), False, PP_ALIGN.CENTER, FONT_LIGHT)

# Bottom insight
add_shape(s, Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.5), GOLD, 0.2)
add_txt(s, Inches(1), Inches(6.82), Inches(11.3), Inches(0.45),
    'المقارنة بين هذه الأركان الثلاثة هي أساس تحليل أداء المشروع',
    14, WHITE, True, PP_ALIGN.CENTER, FONT)

# =====================================================================
# SLIDE 13: CV & SV
# =====================================================================
s = make_slide()
add_shape(s, 0, 0, SW, SH, WHITE)
add_shape(s, 0, 0, Inches(0.12), SH, GOLD)
slide_section(s, 'Cost Variance (CV) & Schedule Variance (SV)', 'انحراف التكلفة وانحراف الجدول')

# CV Box
add_shape(s, Inches(0.8), Inches(1.5), Inches(5.8), Inches(2.4), NAVY)
add_rich(s, Inches(1.2), Inches(1.7), Inches(5.2), Inches(2.0), [
    ('CV  انحراف التكلفة  Cost Variance', 18, GOLD, True, FONT),
    ('', 6, WHITE, False, FONT),
    ('CV = EV  -  AC', 22, WHITE, True, FONT),
    ('', 5, WHITE, False, FONT),
    ('CV > 0  تحت الميزانية  |  CV < 0  تجاوز الميزانية', 13, LIGHT_GRAY, False, FONT),
])
# SV Box
add_shape(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(2.4), NAVY)
add_rich(s, Inches(7.2), Inches(1.7), Inches(5.2), Inches(2.0), [
    ('SV  انحراف الجدول  Schedule Variance', 18, GOLD, True, FONT),
    ('', 6, WHITE, False, FONT),
    ('SV = EV  -  PV', 22, WHITE, True, FONT),
    ('', 5, WHITE, False, FONT),
    ('SV > 0  متقدم عن الجدول  |  SV < 0  متأخر عن الجدول', 13, LIGHT_GRAY, False, FONT),
])

# Example box
add_shape(s, Inches(0.8), Inches(4.2), Inches(11.7), Inches(3), RGBColor(245, 245, 250))
add_txt(s, Inches(1.2), Inches(4.4), Inches(11), Inches(0.4), 'مثال تطبيقي  |  بعد 8 أشهر من مشروع بقيمة 100 مليون ريال', 16, NAVY, True, PP_ALIGN.RIGHT, FONT)
cv_example = [
    ('PV = 78', 'مليون ريال', '(الميزانية المخططة للعمل المجدول)'),
    ('EV = 75', 'مليون ريال', '(قيمة العمل المنجز فعلياً)'),
    ('AC = 82', 'مليون ريال', '(التكاليف الفعلية للعمل المنجز)'),
]
for i, (val, unit, desc) in enumerate(cv_example):
    x = Inches(1.2) + Inches(i * 3.8)
    add_shape(s, x, Inches(5.0), Inches(3.4), Inches(0.9), NAVY, 0.08)
    add_txt(s, x + Inches(0.15), Inches(5.05), Inches(3.1), Inches(0.35), val, 22, GOLD, True, PP_ALIGN.CENTER)
    add_txt(s, x + Inches(0.15), Inches(5.35), Inches(3.1), Inches(0.25), unit + '  |  ' + desc, 10, LIGHT_GRAY, False, PP_ALIGN.CENTER, FONT_LIGHT)

# Results
add_shape(s, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.9), NAVY)
add_txt(s, Inches(1.2), Inches(6.05), Inches(11), Inches(0.4),
    'CV = 75 - 82 = -7 مليون ريال  (تجاوز في التكلفة)    |    SV = 75 - 78 = -3 مليون ريال  (تأخير في الجدول)',
    16, WHITE, True, PP_ALIGN.CENTER, FONT)

# =====================================================================
# SLIDE 14: CPI & SPI
# =====================================================================
s = make_slide()
add_shape(s, 0, 0, SW, SH, WHITE)
add_shape(s, 0, 0, Inches(0.12), SH, GOLD)
slide_section(s, 'Cost Performance Index (CPI) & Schedule Performance Index (SPI)', 'مؤشر أداء التكلفة والجدول')

# CPI Box
add_shape(s, Inches(0.8), Inches(1.5), Inches(5.8), Inches(2.8), NAVY)
add_rich(s, Inches(1.2), Inches(1.7), Inches(5.2), Inches(2.4), [
    ('CPI  مؤشر أداء التكلفة', 18, GOLD, True, FONT),
    ('', 5, WHITE, False, FONT),
    ('CPI = EV / AC', 24, WHITE, True, FONT),
    ('', 5, WHITE, False, FONT),
    ('CPI > 1.0  أداء أفضل من المخطط', 13, TEAL, True, FONT),
    ('CPI = 1.0  أداء مطابق للمخطط', 13, LIGHT_GRAY, False, FONT),
    ('CPI < 1.0  أداء أسوأ من المخطط', 13, RED, True, FONT),
    ('', 4, WHITE, False, FONT),
    ('مثال: CPI = 75/82 = 0.91  (تجاوز في التكلفة)', 14, GOLD, True, FONT),
])
# SPI Box
add_shape(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(2.8), NAVY)
add_rich(s, Inches(7.2), Inches(1.7), Inches(5.2), Inches(2.4), [
    ('SPI  مؤشر أداء الجدول', 18, GOLD, True, FONT),
    ('', 5, WHITE, False, FONT),
    ('SPI = EV / PV', 24, WHITE, True, FONT),
    ('', 5, WHITE, False, FONT),
    ('SPI > 1.0  متقدم عن الجدول', 13, TEAL, True, FONT),
    ('SPI = 1.0  مطابق للجدول', 13, LIGHT_GRAY, False, FONT),
    ('SPI < 1.0  متأخر عن الجدول', 13, RED, True, FONT),
    ('', 4, WHITE, False, FONT),
    ('مثال: SPI = 75/78 = 0.96  (تأخير في الجدول)', 14, GOLD, True, FONT),
])

# Interpretation table
add_shape(s, Inches(0.8), Inches(4.6), Inches(11.7), Inches(2.6), RGBColor(245, 245, 250))
add_txt(s, Inches(1.2), Inches(4.7), Inches(11), Inches(0.35), 'دليل تفسير مؤشرات الأداء  |  Performance Interpretation Guide', 15, NAVY, True, PP_ALIGN.CENTER)

headers = ['الحالة / Status', 'CPI', 'SPI', 'التفسير / Interpretation']
rows = [
    ['ممتاز / Excellent', '> 1.0', '> 1.0', 'تحت الميزانية ومتقدم عن الجدول'],
    ['جيد / Good', '> 1.0', '< 1.0', 'تحت الميزانية ولكن متأخر'],
    ['تحت المراقبة / Watch', '< 1.0', '> 1.0', 'تجاوز في التكلفة ولكن متقدم'],
    ['تحتاج تدخل / Critical', '< 1.0', '< 1.0', 'تجاوز في التكلفة ومتأخر'],
]
for j, row in enumerate(rows):
    y = Inches(5.1) + Inches(j * 0.38)
    for k, cell in enumerate(row):
        x = Inches(1.2) + Inches(k * 2.9)
        c = NAVY if j == 0 else RGBColor(60, 60, 60)
        b = j == 0
        add_txt(s, x, y, Inches(2.8), Inches(0.35), cell, 10 if k < 3 else 11, c, b, PP_ALIGN.CENTER, FONT_LIGHT)

# =====================================================================
# SLIDE 15: CPI/SPI CHART
# =====================================================================
s = make_slide()
add_full_bg(s, chart('02_cpi_spi_trend.png'))
# Title bar
add_shape(s, 0, Inches(6.3), SW, Inches(1.2), NAVY)
add_shape(s, 0, Inches(6.3), SW, Pt(3), GOLD)
add_txt(s, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.8),
    'اتجاه مؤشرات الأداء CPI و SPI  |  تحليل الأداء الشهري',
    18, WHITE, False, PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 16: EAC & ETC
# =====================================================================
s = make_slide()
add_shape(s, 0, 0, SW, SH, WHITE)
add_shape(s, 0, 0, Inches(0.12), SH, GOLD)
slide_section(s, 'EAC & ETC  |  Estimate at Completion & Estimate to Complete', 'تقدير التكلفة عند الإكمال والمتبقي')

# EAC Box
add_shape(s, Inches(0.8), Inches(1.5), Inches(5.8), Inches(2.6), NAVY)
add_rich(s, Inches(1.2), Inches(1.7), Inches(5.2), Inches(2.2), [
    ('EAC  تقدير التكلفة عند الإكمال', 17, GOLD, True, FONT),
    ('', 4, WHITE, False, FONT),
    ('EAC = BAC / CPI', 22, WHITE, True, FONT),
    ('', 4, WHITE, False, FONT),
    ('أو:  EAC = AC + (BAC - EV) / CPI', 15, LIGHT_GRAY, False, FONT),
    ('', 4, WHITE, False, FONT),
    ('مثال:  EAC = 100 / 0.91 = 109.9 مليون ريال', 13, GOLD, True, FONT),
    ('نتوقع أن تكلفة المشروع النهائية 109.9 مليون', 12, LIGHT_GRAY, False, FONT),
])
# ETC Box
add_shape(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(2.6), NAVY)
add_rich(s, Inches(7.2), Inches(1.7), Inches(5.2), Inches(2.2), [
    ('ETC  تقدير التكلفة المتبقية', 17, GOLD, True, FONT),
    ('', 4, WHITE, False, FONT),
    ('ETC = EAC  -  AC', 22, WHITE, True, FONT),
    ('', 8, WHITE, False, FONT),
    ('مثال:  ETC = 109.9 - 82 = 27.9 مليون ريال', 13, GOLD, True, FONT),
    ('التكلفة المتبقية لإكمال المشروع', 12, LIGHT_GRAY, False, FONT),
])

# EAC Formulas detail
add_shape(s, Inches(0.8), Inches(4.4), Inches(11.7), Inches(2.8), RGBColor(245, 245, 250))
add_txt(s, Inches(1.2), Inches(4.5), Inches(11), Inches(0.35), 'طرق حساب EAC  |  EAC Calculation Methods', 15, NAVY, True, PP_ALIGN.CENTER)
eac_methods = [
    ('EAC = BAC / CPI', 'يفترض استمرار أداء التكلفة الحالي', 'Typical performance continues'),
    ('EAC = AC + (BAC - EV)', 'يفترض أداء مطابق للمخطط للفترة المتبقية', 'Remaining work at budgeted rate'),
    ('EAC = AC + (BAC - EV) / (CPI × SPI)', 'يفترض تأثر التكلفة والجدول معاً', 'Cost & schedule factors combined'),
    ('EAC = AC + ETC (جديد)', 'تقدير جديد من الموقع للمتبقي', 'New bottom-up estimate'),
]
for i, (formula, ar_desc, en_desc) in enumerate(eac_methods):
    y = Inches(4.9) + Inches(i * 0.52)
    add_shape(s, Inches(1.2), y, Inches(2.5), Inches(0.42), TEAL if i % 2 == 0 else GOLD, 0.15)
    add_txt(s, Inches(1.3), y + Pt(2), Inches(2.3), Inches(0.35), formula, 11, TEAL if i % 2 == 0 else GOLD, True, PP_ALIGN.CENTER, FONT)
    add_txt(s, Inches(3.8), y + Pt(2), Inches(4.0), Inches(0.35), ar_desc, 11, RGBColor(60, 60, 60), False, PP_ALIGN.RIGHT)
    add_txt(s, Inches(7.8), y + Pt(2), Inches(4.5), Inches(0.35), f'({en_desc})', 10, MED_GRAY, False, PP_ALIGN.LEFT, FONT_LIGHT)

# =====================================================================
# SLIDE 17: VAC & TCPI
# =====================================================================
s = make_slide()
add_shape(s, 0, 0, SW, SH, WHITE)
add_shape(s, 0, 0, Inches(0.12), SH, GOLD)
slide_section(s, 'VAC & TCPI  |  Variance at Completion & To-Complete Performance Index', 'انحراف الإكمال ومؤشر الأداء المطلوب')

add_shape(s, Inches(0.8), Inches(1.5), Inches(5.8), Inches(2.4), NAVY)
add_rich(s, Inches(1.2), Inches(1.7), Inches(5.2), Inches(2.0), [
    ('VAC  انحراف التكلفة عند الإكمال', 17, GOLD, True, FONT),
    ('', 5, WHITE, False, FONT),
    ('VAC = BAC  -  EAC', 22, WHITE, True, FONT),
    ('', 5, WHITE, False, FONT),
    ('VAC > 0  تحت الميزانية المتوقعة', 13, TEAL, True, FONT),
    ('VAC < 0  تجاوز في الميزانية المتوقعة', 13, RED, True, FONT),
    ('', 4, WHITE, False, FONT),
    ('مثال:  VAC = 100 - 109.9 = -9.9 مليون', 14, GOLD, True, FONT),
    ('نتوقع تجاوز الميزانية بـ 9.9 مليون ريال', 12, LIGHT_GRAY, False, FONT),
])
add_shape(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(2.4), NAVY)
add_rich(s, Inches(7.2), Inches(1.7), Inches(5.2), Inches(2.0), [
    ('TCPI  مؤشر الأداء المطلوب لإكمال', 17, GOLD, True, FONT),
    ('', 5, WHITE, False, FONT),
    ('TCPI = (BAC - EV) / (BAC - AC)', 18, WHITE, True, FONT),
    ('أو  TCPI = (BAC - EV) / (EAC - AC)', 15, LIGHT_GRAY, False, FONT),
    ('', 5, WHITE, False, FONT),
    ('TCPI > 1.0  أداء أكثر كفاءة مطلوب', 13, RED, True, FONT),
    ('TCPI < 1.0  أداء أقل صرامة ممكن', 13, TEAL, True, FONT),
    ('', 4, WHITE, False, FONT),
    ('مثال:  TCPI = (100-75)/(100-82) = 1.39', 13, GOLD, True, FONT),
])

# Comparison box
add_shape(s, Inches(0.8), Inches(4.2), Inches(11.7), Inches(3), RGBColor(245, 245, 250))
add_txt(s, Inches(1.2), Inches(4.3), Inches(11), Inches(0.35), 'ملخص مؤشرات EVM  |  EVM Indicators Summary', 16, NAVY, True, PP_ALIGN.CENTER)
summary_data = [
    ('BAC', '100.0', 'الميزانية المعتمدة', GREEN),
    ('EAC', '109.9', 'تقدير التكلفة النهائية', RED),
    ('ETC', '27.9', 'التكلفة المتبقية للإكمال', TEAL),
    ('VAC', '-9.9', 'الانحراف المتوقع عند الإكمال', RED),
    ('TCPI', '1.39', 'مؤشر الأداء المطلوب', RED),
]
for i, (ind, val, desc, color) in enumerate(summary_data):
    x = Inches(1.0) + Inches(i * 2.35)
    add_shape(s, x, Inches(4.8), Inches(2.15), Inches(2.1), NAVY)
    add_shape(s, x, Inches(4.8), Inches(2.15), Inches(0.55), color)
    add_txt(s, x, Inches(4.85), Inches(2.15), Inches(0.45), ind, 18, WHITE, True, PP_ALIGN.CENTER)
    add_txt(s, x, Inches(5.5), Inches(2.15), Inches(0.45), val, 24, WHITE, True, PP_ALIGN.CENTER)
    add_txt(s, x, Inches(6.1), Inches(2.15), Inches(0.5), desc, 9, LIGHT_GRAY, False, PP_ALIGN.CENTER, FONT_LIGHT)

# =====================================================================
# SLIDE 18: FORMULA SUMMARY
# =====================================================================
s = make_slide()
add_shape(s, 0, 0, SW, SH, NAVY)
add_shape(s, 0, 0, Inches(0.12), SH, GOLD)
slide_section(s, 'EVM Formula Summary Table', 'جدول ملخص معادلات EVM')

formulas = [
    ('PV', 'Planned Value', 'BAC × % مخطط', 'القيمة المخططة'),
    ('EV', 'Earned Value', 'BAC × % منجز', 'القيمة المكتسبة'),
    ('AC', 'Actual Cost', 'مجموع التكاليف', 'التكلفة الفعلية'),
    ('CV', 'Cost Variance', 'EV - AC', 'انحراف التكلفة'),
    ('SV', 'Schedule Variance', 'EV - PV', 'انحراف الجدول'),
    ('CPI', 'Cost Perf. Index', 'EV / AC', 'مؤشر أداء التكلفة'),
    ('SPI', 'Schedule Perf. Index', 'EV / PV', 'مؤشر أداء الجدول'),
    ('EAC', 'Est. at Completion', 'BAC / CPI', 'تقدير التكلفة النهائية'),
    ('ETC', 'Est. to Complete', 'EAC - AC', 'تقدير التكلفة المتبقية'),
    ('VAC', 'Variance at Compl.', 'BAC - EAC', 'انحراف التكلفة النهائي'),
    ('TCPI', 'To-Complete Perf. Index', '(BAC-EV)/(BAC-AC)', 'مؤشر الأداء المطلوب'),
]
y_start = Inches(1.6)
# Table header
add_shape(s, Inches(0.6), y_start, Inches(12.1), Inches(0.45), GOLD)
for j, (hdr, w) in enumerate([('المؤشر', 1.2), ('English', 2.0), ('Formula / المعادلة', 4.5), ('الوصف', 4.4)]):
    x = Inches(0.8) + sum([1.2, 2.0, 4.5, 4.4][:j]) if False else \
        [Inches(0.8), Inches(2.1), Inches(4.2), Inches(8.8)][j]
    add_txt(s, x, y_start + Pt(2), Inches([1.2, 2.0, 4.5, 4.4][j]), Inches(0.35),
            hdr, 12, NAVY, True, PP_ALIGN.CENTER if j in (0, 2) else PP_ALIGN.RIGHT if j == 3 else PP_ALIGN.LEFT)

for i, (ind, en, formula, desc) in enumerate(formulas):
    y = y_start + Inches(0.5) + Inches(i * 0.42)
    bg = RGBColor(35, 50, 80) if i % 2 == 0 else RGBColor(42, 58, 90)
    add_shape(s, Inches(0.6), y, Inches(12.1), Inches(0.4), bg)
    add_txt(s, Inches(0.8), y + Pt(3), Inches(1.2), Inches(0.3), ind, 14, GOLD, True, PP_ALIGN.CENTER)
    add_txt(s, Inches(2.1), y + Pt(3), Inches(2.0), Inches(0.3), en, 11, LIGHT_GRAY, False, PP_ALIGN.LEFT, FONT_LIGHT)
    add_txt(s, Inches(4.2), y + Pt(3), Inches(4.5), Inches(0.3), formula, 13, WHITE, True, PP_ALIGN.CENTER)
    add_txt(s, Inches(8.8), y + Pt(3), Inches(4.4), Inches(0.3), desc, 11, LIGHT_GRAY, False, PP_ALIGN.RIGHT)

# =====================================================================
# SLIDE 19: S-CURVE
# =====================================================================
s = make_slide()
add_full_bg(s, chart('01_scurve.png'))
add_shape(s, 0, Inches(6.3), SW, Inches(1.2), NAVY)
add_shape(s, 0, Inches(6.3), SW, Pt(3), GOLD)
add_txt(s, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.8),
    'منحنى S  |  مقارنة القيمة المخططة والمكتسبة والتكلفة الفعلية',
    18, WHITE, False, PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 20: CV/SV CHART
# =====================================================================
s = make_slide()
add_full_bg(s, chart('03_cv_sv.png'))
add_shape(s, 0, Inches(6.3), SW, Inches(1.2), NAVY)
add_shape(s, 0, Inches(6.3), SW, Pt(3), GOLD)
add_txt(s, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.8),
    'انحراف التكلفة والجدول  |  تحليل شهري',
    18, WHITE, False, PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 21: EAC FORECAST CHART
# =====================================================================
s = make_slide()
add_full_bg(s, chart('04_eac_forecast.png'))
add_shape(s, 0, Inches(6.3), SW, Inches(1.2), NAVY)
add_shape(s, 0, Inches(6.3), SW, Pt(3), GOLD)
add_txt(s, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.8),
    'تقدير التكلفة عند الإكمال (EAC)  |  توقعات نحو BAC',
    18, WHITE, False, PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 22: VAC TREND CHART
# =====================================================================
s = make_slide()
add_full_bg(s, chart('08_vac_trend.png'))
add_shape(s, 0, Inches(6.3), SW, Inches(1.2), NAVY)
add_shape(s, 0, Inches(6.3), SW, Pt(3), GOLD)
add_txt(s, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.8),
    'انحراف الإكمال (VAC)  |  التعافي نحو الهدف',
    18, WHITE, False, PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 23: EAC/ETC/VAC COMPARISON CHART
# =====================================================================
s = make_slide()
add_full_bg(s, chart('06_eac_etc_vac.png'))
add_shape(s, 0, Inches(6.3), SW, Inches(1.2), NAVY)
add_shape(s, 0, Inches(6.3), SW, Pt(3), GOLD)
add_txt(s, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.8),
    'مقارنة BAC / EAC / ETC / VAC  |  ملخص توقعات المشروع',
    18, WHITE, False, PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 24: COST BREAKDOWN CHART
# =====================================================================
s = make_slide()
add_full_bg(s, chart('07_cost_breakdown.png'))
add_shape(s, 0, Inches(6.3), SW, Inches(1.2), NAVY)
add_shape(s, 0, Inches(6.3), SW, Pt(3), GOLD)
add_txt(s, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.8),
    'هيكل انهيار التكاليف في المشاريع الإنشائية',
    18, WHITE, False, PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 25: PERFORMANCE DASHBOARD
# =====================================================================
s = make_slide()
add_full_bg(s, chart('05_dashboard_gauges.png'))
add_shape(s, 0, Inches(6.3), SW, Inches(1.2), NAVY)
add_shape(s, 0, Inches(6.3), SW, Pt(3), GOLD)
add_txt(s, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.8),
    'لوحة قيادة أداء EVM  |  مؤشرات الفترة الحالية',
    18, WHITE, False, PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 26: EVM IN PLANNING
# =====================================================================
s = make_slide()
add_shape(s, 0, 0, SW, SH, WHITE)
add_shape(s, 0, 0, Inches(0.12), SH, GOLD)
slide_section(s, 'EVM in the Planning Phase', 'تطبيق EVM في مرحلة التخطيط')

add_shape(s, Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.2), NAVY)
add_rich(s, Inches(1.2), Inches(1.7), Inches(5.2), Inches(4.8), [
    ('مرحلة التخطيط', 20, GOLD, True, FONT),
    ('', 8, WHITE, False, FONT),
    ('1. إنشاء WBS', 14, TEAL, True, FONT),
    ('تقسيم المشروع إلى حزم عمل قابلة للقياس', 12, LIGHT_GRAY, False, FONT),
    ('', 5, WHITE, False, FONT),
    ('2. تخصيص الميزانية', 14, TEAL, True, FONT),
    ('توزيع BAC على كل حزمة عمل', 12, LIGHT_GRAY, False, FONT),
    ('', 5, WHITE, False, FONT),
    ('3. إعداد الجدول الزمني', 14, TEAL, True, FONT),
    ('توزيع القيمة المخططة (PV) زمنياً', 12, LIGHT_GRAY, False, FONT),
    ('', 5, WHITE, False, FONT),
    ('4. تحديد نقاط القياس', 14, TEAL, True, FONT),
    ('طرق قياس الإنجاز (Weighted Milestones, % Complete)', 12, LIGHT_GRAY, False, FONT),
    ('', 5, WHITE, False, FONT),
    ('5. اعتماد خط الأساس', 14, GOLD, True, FONT),
    ('Baseline = التوزيع الزمني لـ PV', 12, LIGHT_GRAY, False, FONT),
])
add_shape(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.2), RGBColor(245, 245, 250))
add_txt(s, Inches(7.1), Inches(1.7), Inches(5.3), Inches(0.4), 'Construction Planning Example', 16, NAVY, True, PP_ALIGN.LEFT)
plan_items = [
    'مشروع بناء برج سكني 20 دور',
    'BAC = 200 مليون ريال',
    'المدة: 24 شهراً',
    '',
    'WBS Level 1: أعمال هيكلية',
    '  - أساسات: 30 مليون ريال (مدة 4 أشهر)',
    '  - هيكل خرساني: 50 مليون ريال (مدة 8 أشهر)',
    '  - واجهات: 20 مليون ريال (مدة 6 أشهر)',
    'WBS Level 2: أعمال تشطيب',
    '  - كهرباء: 25 مليون ريال (مدة 6 أشهر)',
    '  - سباكة: 15 مليون ريال (مدة 5 أشهر)',
    '  - تكييف: 20 مليون ريال (مدة 5 أشهر)',
    'WBS Level 3: أعمال خارجية',
    '  - ساحات: 20 مليون ريال (مدة 4 أشهر)',
    '  - خدمات: 20 مليون ريال (مدة 6 أشهر)',
]
for i, item in enumerate(plan_items):
    c = NAVY if item.endswith(':') or 'WBS' in item else RGBColor(80, 80, 80)
    b = item.endswith(':') or 'WBS' in item
    add_txt(s, Inches(7.1), Inches(2.2) + Inches(i * 0.24), Inches(5.3), Inches(0.22), item, 9, c, b, PP_ALIGN.LEFT, FONT_LIGHT)

# =====================================================================
# SLIDE 27: EVM IN EXECUTION
# =====================================================================
s = make_slide()
add_full_bg(s, img('01_Wadi_Cascade_Waterfall_1920x1080-Gross.jpg'))
add_shape(s, 0, 0, SW, SH, DARK_BG, 0.6)
slide_title_bar(s, 'EVM in Execution & Monitoring', 'تطبيق EVM في مرحلة التنفيذ والمراقبة')

processes = [
    ('تسجيل التكاليف\nالفعلية', 'Track actual\ncosts from ERP', TEAL),
    ('قياس\nالإنجاز', 'Measure physical\n% complete', GOLD),
    ('حساب\nالمؤشرات', 'Calculate CPI,\nSPI, CV, SV', NAVY),
    ('تحليل\nالانحرافات', 'Analyze variances\n& root causes', TEAL),
    ('التنبؤ\nبالنتائج', 'Forecast EAC,\nETC, VAC', GOLD),
    ('إعداد\nالتقارير', 'Generate PMO\ndashboard reports', NAVY),
]
for i, (ar, en, color) in enumerate(processes):
    x = Inches(0.5) + Inches(i * 2.1)
    add_shape(s, x, Inches(1.6), Inches(1.9), Inches(2.2), color, 0.2)
    add_shape(s, x, Inches(1.6), Inches(1.9), Inches(1.0), color, 0.85)
    add_txt(s, x, Inches(1.65), Inches(1.9), Inches(0.9), ar, 13, WHITE, True, PP_ALIGN.CENTER)
    add_txt(s, x, Inches(2.7), Inches(1.9), Inches(0.8), en, 9, RGBColor(230, 230, 230), False, PP_ALIGN.CENTER, FONT_LIGHT)
    if i < len(processes) - 1:
        add_shape(s, x + Inches(1.9), Inches(2.5), Inches(0.2), Pt(2), GOLD)

# Bottom section
add_shape(s, Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.8), NAVY, 0.85)
add_txt(s, Inches(0.8), Inches(4.4), Inches(11.5), Inches(0.35), 'Reporting Cycle for Construction Projects', 15, GOLD, True, PP_ALIGN.CENTER)
report_table = [
    ('التقرير / Report', 'الدورية', 'المحتوى', 'الجمهور'),
    ('Daily Progress', 'يومي', 'الإنجاز اليومي، المشاكل', 'موقع المشروع'),
    ('Weekly EVM', 'أسبوعي', 'CPI, SPI, CV, SV', 'إدارة المشروع'),
    ('Monthly PMO', 'شهري', 'EAC, ETC, VAC, Forecast', 'PMO / الإدارة العليا'),
    ('Quarterly Board', 'ربع سنوي', 'Overall Dashboard', 'مجلس الإدارة / العميل'),
]
for j, row in enumerate(report_table):
    y = Inches(4.85) + Inches(j * 0.4)
    bg = GOLD if j == 0 else DARK_NAVY
    c = NAVY if j == 0 else WHITE
    b = j == 0
    add_shape(s, Inches(0.8), y, Inches(11.7), Inches(0.38), bg)
    for k, cell in enumerate(row):
        xs = [Inches(0.9), Inches(4.0), Inches(6.5), Inches(9.0)][k]
        ws = [3.0, 2.4, 2.4, 2.4][k]
        sz = 11 if j == 0 else 10
        add_txt(s, xs, y + Pt(2), ws, Inches(0.32), cell, sz, c, b, PP_ALIGN.CENTER, FONT_LIGHT)

# =====================================================================
# SLIDE 28: EVM IN MEGA PROJECTS
# =====================================================================
s = make_slide()
add_full_bg(s, img('ksp_masterplan.jpg'))
add_shape(s, 0, 0, SW, SH, DARK_BG, 0.6)
slide_title_bar(s, 'EVM in Mega Projects', 'تطبيق EVM في المشاريع العملاقة')

add_rich(s, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.5), [
    ('تحديات المشاريع العملاقة', 18, GOLD, True, FONT),
    ('', 8, WHITE, False, FONT),
    ('•  آلاف حزم العمل (WBS)', 13, WHITE, False, FONT),
    ('•  مقاولون متعددون ومتداخلون', 13, WHITE, False, FONT),
    ('•  فترات زمنية طويلة (5-10 سنوات)', 13, WHITE, False, FONT),
    ('•  ميزانيات ضخمة (مليارات)', 13, WHITE, False, FONT),
    ('•  تغيرات في النطاق', 13, WHITE, False, FONT),
    ('•  مخاطر تراكمية', 13, WHITE, False, FONT),
    ('', 8, WHITE, False, FONT),
    ('EVM هو الحل الأمثل لـ:', 15, GOLD, True, FONT),
    ('•  توحيد مقاييس الأداء', 13, WHITE, False, FONT),
    ('•  ربط التكاليف بالجدول', 13, WHITE, False, FONT),
    ('•  الشفافية مع أصحاب المصلحة', 13, WHITE, False, FONT),
    ('•  التنبؤ المبكر', 13, WHITE, False, FONT),
], PP_ALIGN.RIGHT)

add_shape(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.5), NAVY, 0.85)
add_txt(s, Inches(7.1), Inches(1.7), Inches(5.3), Inches(0.35), 'Mega Project EVM Implementation', 15, GOLD, True, PP_ALIGN.LEFT)
mega_items = [
    'Neom City Projects',
    '  • Multiple packages, single EVM framework',
    '  • Integrated with SAP and Primavera P6',
    '',
    'King Salman Park, Riyadh',
    '  • 16.6 km², SAR 23 billion budget',
    '  • Roll-up EVM from contractors to PMO',
    '',
    'Riyadh Metro',
    '  • 6 lines, 176 km, SAR 85 billion',
    '  • EVM at program level with cascading KPIs',
    '',
    'Key Success Factors:',
    '  ■ Centralized PMIS / ERP integration',
    '  ■ Standardized WBS across all packages',
    '  ■ Automated data collection',
    '  ■ Monthly EVM governance reviews',
]
for i, item in enumerate(mega_items):
    c = GOLD if item.startswith('  ') or item.startswith('■') else NAVY if item else MED_GRAY
    b = not item.startswith('  ') and bool(item) and not item.startswith('■')
    if item == '':
        continue
    add_txt(s, Inches(7.1), Inches(2.2) + Inches(i * 0.24), Inches(5.3), Inches(0.22), item, 9, c if item else MED_GRAY, b, PP_ALIGN.LEFT, FONT_LIGHT)

# =====================================================================
# SLIDE 29: CASE STUDY INTRO
# =====================================================================
s = make_slide()
add_shape(s, 0, 0, SW, SH, WHITE)
add_shape(s, 0, 0, Inches(0.12), SH, GOLD)
slide_section(s, 'Case Study  |  Real-World EVM Application', 'دراسة حالة عملية')

add_shape(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.2), NAVY, 0.05)
add_txt(s, Inches(1.2), Inches(1.7), Inches(11), Inches(0.4), 'مشروع إنشاء مجمع تجاري وسكني  |  Commercial & Residential Complex', 18, GOLD, True, PP_ALIGN.RIGHT)

# Project details
details_left = [
    ('مواصفات المشروع / Project Specifications', True, 14, GOLD),
    ('BAC  =  100,000,000 ريال سعودي', False, 13, WHITE),
    ('المدة  =  12 شهراً', False, 13, WHITE),
    ('الموقع  =  الرياض، المملكة العربية السعودية', False, 13, WHITE),
    ('المساحة  =  25,000 م²', False, 13, WHITE),
    ('طبيعة العمل  =  أعمال هيكلية + تشطيب + بنية تحتية', False, 13, WHITE),
    ('طريقة الدفع  =  دفعات مرحلية حسب الإنجاز', False, 13, WHITE),
]
for i, (txt, b, sz, c) in enumerate(details_left):
    add_txt(s, Inches(1.2), Inches(2.2) + Inches(i * 0.38), Inches(5.5), Inches(0.35), txt, sz, c, b, PP_ALIGN.RIGHT, FONT_LIGHT if not b else FONT)

details_right = [
    ('حالة المشروع بعد 8 أشهر / Status at Month 8', True, 14, TEAL),
    ('النسبة المخطط إنجازها  =  78%', False, 13, WHITE),
    ('PV  =  78,000,000 ريال', False, 13, WHITE),
    ('النسبة المنجزة فعلياً  =  75%', False, 13, WHITE),
    ('EV  =  75,000,000 ريال', False, 13, WHITE),
    ('التكاليف الفعلية  =  82,000,000 ريال', False, 13, WHITE),
    ('AC  =  82,000,000 ريال', False, 13, WHITE),
]
for i, (txt, b, sz, c) in enumerate(details_right):
    add_txt(s, Inches(7), Inches(2.2) + Inches(i * 0.38), Inches(5.5), Inches(0.35), txt, sz, c, b, PP_ALIGN.LEFT, FONT_LIGHT if not b else FONT)

# Divider line
add_shape(s, Inches(6.5), Inches(2.0), Pt(2), Inches(4.5), GOLD)

# Bottom
add_shape(s, Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.7), GOLD, 0.15)
add_txt(s, Inches(1.2), Inches(5.85), Inches(11), Inches(0.5),
    'سنقوم الآن بحساب وتحليل جميع مؤشرات EVM لهذا المشروع',
    16, WHITE, True, PP_ALIGN.CENTER, FONT)

# =====================================================================
# SLIDE 30: CASE STUDY CALCULATIONS
# =====================================================================
s = make_slide()
add_shape(s, 0, 0, SW, SH, NAVY)
add_shape(s, 0, 0, Inches(0.12), SH, GOLD)
slide_section(s, 'Case Study  |  EVM Calculations', 'حسابات القيمة المكتسبة')

calc_items = [
    ('CV', 'EV - AC', '75 - 82', '-7 مليون ريال', RED, 'تجاوز في التكلفة'),
    ('SV', 'EV - PV', '75 - 78', '-3 مليون ريال', RED, 'تأخير في الجدول'),
    ('CPI', 'EV / AC', '75 / 82', '0.91', GOLD, 'أقل من 1 (غير مرضٍ)'),
    ('SPI', 'EV / PV', '75 / 78', '0.96', GOLD, 'أقل من 1 (متأخر)'),
    ('EAC', 'BAC / CPI', '100 / 0.91', '109.9 مليون', RED, 'متوقع تجاوز الميزانية'),
    ('ETC', 'EAC - AC', '109.9 - 82', '27.9 مليون', TEAL, 'التكلفة المتبقية'),
    ('VAC', 'BAC - EAC', '100 - 109.9', '-9.9 مليون', RED, 'انحراف نهائي متوقع'),
    ('TCPI', '(BAC-EV)/(BAC-AC)', '25 / 18', '1.39', RED, 'أداء أعلى مطلوب'),
]
for i, (ind, formula, calc, result, color, desc) in enumerate(calc_items):
    y = Inches(1.5) + Inches(i * 0.7)
    add_shape(s, Inches(0.8), y, Inches(11.7), Inches(0.62), RGBColor(30, 45, 72) if i % 2 == 0 else RGBColor(36, 52, 80))
    add_shape(s, Inches(0.8), y, Inches(1.0), Inches(0.62), color)
    add_txt(s, Inches(0.85), y + Pt(4), Inches(0.9), Inches(0.5), ind, 18, WHITE, True, PP_ALIGN.CENTER)
    add_txt(s, Inches(2.0), y + Pt(4), Inches(2.5), Inches(0.5), formula, 13, LIGHT_GRAY, False, PP_ALIGN.CENTER, FONT_LIGHT)
    add_txt(s, Inches(4.7), y + Pt(4), Inches(2.0), Inches(0.5), calc, 13, LIGHT_GRAY, False, PP_ALIGN.CENTER, FONT_LIGHT)
    add_txt(s, Inches(6.9), y + Pt(4), Inches(2.5), Inches(0.5), result, 16, color, True, PP_ALIGN.CENTER)
    add_txt(s, Inches(9.6), y + Pt(4), Inches(2.7), Inches(0.5), desc, 11, LIGHT_GRAY, False, PP_ALIGN.CENTER, FONT_LIGHT)

# =====================================================================
# SLIDE 31: CASE STUDY ANALYSIS
# =====================================================================
s = make_slide()
add_shape(s, 0, 0, SW, SH, WHITE)
add_shape(s, 0, 0, Inches(0.12), SH, GOLD)
slide_section(s, 'Case Study  |  Analysis & Recommendations', 'تحليل وتوصيات')

add_shape(s, Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.2), NAVY)
add_rich(s, Inches(1.2), Inches(1.7), Inches(5.2), Inches(4.8), [
    ('تحليل النتائج', 18, GOLD, True, FONT),
    ('', 6, WHITE, False, FONT),
    ('تكلفة:', 14, TEAL, True, FONT),
    ('المشروع تجاوز الميزانية بنسبة 9%', 13, WHITE, False, FONT),
    ('CPI = 0.91 أقل من المستهدف', 13, WHITE, False, FONT),
    ('السبب: ارتفاع أسعار المواد', 13, LIGHT_GRAY, False, FONT),
    ('', 5, WHITE, False, FONT),
    ('جدول:', 14, TEAL, True, FONT),
    ('المشروع متأخر بنسبة 4%', 13, WHITE, False, FONT),
    ('SPI = 0.96 أقل من المطلوب', 13, WHITE, False, FONT),
    ('السبب: تأخر المقاول من الباطن', 13, LIGHT_GRAY, False, FONT),
    ('', 5, WHITE, False, FONT),
    ('توقعات:', 14, GOLD, True, FONT),
    ('EAC = 109.9 مليون ريال', 13, WHITE, False, FONT),
    ('TCPI = 1.39 أداء أعلى مطلوب', 13, WHITE, False, FONT),
])
add_shape(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.2), RGBColor(245, 245, 250))
add_txt(s, Inches(7.1), Inches(1.7), Inches(5.3), Inches(0.4), 'التوصيات  |  Recommendations', 16, NAVY, True, PP_ALIGN.RIGHT)

recommendations = [
    ('1. تسريع الأعمال الحرجة', 'Accelerate critical path activities',
     'زيادة الورديات أو تعيين مقاول إضافي للأنشطة الحرجة'),
    ('2. تحسين الإنتاجية', 'Improve productivity',
     'مراجعة أداء المقاولين وتحسين كفاءة العمل'),
    ('3. مراقبة التكاليف', 'Cost control measures',
     'تشديد الرقابة على المشتريات والعقود'),
    ('4. إدارة التغيير', 'Change management',
     'مراجعة طلبات التغيير وتأثيرها على الميزانية'),
    ('5. تحديث التوقعات', 'Update forecasts',
     'تحديث EAC شهرياً ليعكس الواقع'),
]
for i, (ar, en, detail) in enumerate(recommendations):
    y = Inches(2.3) + Inches(i * 0.9)
    add_shape(s, Inches(7.1), y, Inches(5.3), Inches(0.8), NAVY, 0.05)
    add_txt(s, Inches(7.3), y + Pt(3), Inches(4.8), Inches(0.3), ar, 14, GOLD, True, PP_ALIGN.RIGHT)
    add_txt(s, Inches(7.3), y + Pt(18), Inches(4.8), Inches(0.25), f'({en})', 10, TEAL, False, PP_ALIGN.RIGHT, FONT_LIGHT)
    add_txt(s, Inches(7.3), y + Pt(30), Inches(4.8), Inches(0.3), detail, 10, RGBColor(80, 80, 80), False, PP_ALIGN.RIGHT, FONT_LIGHT)

# =====================================================================
# SLIDE 32: EVM WITH PRIMAVERA P6
# =====================================================================
s = make_slide()
add_shape(s, 0, 0, SW, SH, WHITE)
add_shape(s, 0, 0, Inches(0.12), SH, GOLD)
slide_section(s, 'EVM Integration with Primavera P6', 'التكامل مع برنامج Primavera P6')

add_shape(s, Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.2), NAVY)
add_rich(s, Inches(1.2), Inches(1.7), Inches(5.2), Inches(4.8), [
    ('التكامل مع P6', 18, GOLD, True, FONT),
    ('', 6, WHITE, False, FONT),
    ('Primavera P6 هو النظام الأكثر استخداماً', 13, WHITE, False, FONT),
    ('لإدارة الجداول الزمنية في المشاريع', 13, WHITE, False, FONT),
    ('الإنشائية الكبرى.', 13, WHITE, False, FONT),
    ('', 5, WHITE, False, FONT),
    ('ميزات التكامل:', 14, TEAL, True, FONT),
    ('•  ربط WBS مباشرة مع P6', 12, WHITE, False, FONT),
    ('•  توزيع PV تلقائياً حسب الجدول', 12, WHITE, False, FONT),
    ('•  تحديث EV بناءً على % Complete', 12, WHITE, False, FONT),
    ('•  حساب SPI من تقدم الجدول', 12, WHITE, False, FONT),
    ('•  تصدير CPI و SPI إلى PMO', 12, WHITE, False, FONT),
    ('', 5, WHITE, False, FONT),
    ('خطوات التكامل:', 14, GOLD, True, FONT),
    ('1. إنشاء WBS في P6', 12, WHITE, False, FONT),
    ('2. ربط الموارد والتكاليف', 12, WHITE, False, FONT),
    ('3. تعيين Budget و Planned Values', 12, WHITE, False, FONT),
    ('4. تحديث % Complete أسبوعياً', 12, WHITE, False, FONT),
    ('5. حساب EVM تلقائياً', 12, WHITE, False, FONT),
])
add_shape(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.2), RGBColor(245, 245, 250))
add_txt(s, Inches(7.1), Inches(1.7), Inches(5.3), Inches(0.4), 'P6 EVM Fields', 16, NAVY, True, PP_ALIGN.LEFT)
p6_fields = [
    ('P6 Field', 'EVM Equivalent', 'Description'),
    ('Budgeted Total Cost', 'BAC', 'إجمالي الميزانية'),
    ('Budgeted Cost (BCWS)', 'PV', 'القيمة المخططة'),
    ('Earned Value (BCWP)', 'EV', 'القيمة المكتسبة'),
    ('Actual Cost (ACWP)', 'AC', 'التكلفة الفعلية'),
    ('Cost Variance', 'CV', 'انحراف التكلفة'),
    ('Schedule Variance', 'SV', 'انحراف الجدول'),
    ('CPI (Cost Perf. Index)', 'CPI', 'مؤشر أداء التكلفة'),
    ('SPI (Schedule Perf.)', 'SPI', 'مؤشر أداء الجدول'),
    ('EAC (Est. at Compl.)', 'EAC', 'تقدير الإكمال'),
    ('ETC (Est. to Compl.)', 'ETC', 'تقدير المتبقي'),
    ('VAC', 'VAC', 'انحراف الإكمال'),
    ('TCPI', 'TCPI', 'مؤشر الأداء المطلوب'),
]
for j, row in enumerate(p6_fields):
    y = Inches(2.2) + Inches(j * 0.35)
    bg = GOLD if j == 0 else DARK_NAVY if j % 2 == 0 else NAVY
    c = NAVY if j == 0 else WHITE
    b = j == 0
    add_shape(s, Inches(7.1), y, Inches(5.3), Inches(0.33), bg)
    for k, cell in enumerate(row):
        xs = [Inches(7.2), Inches(9.5), Inches(10.5)][k]
        ws = [2.2, 1.0, 1.6][k]
        add_txt(s, xs, y + Pt(1), ws, Inches(0.28), cell, 8, c, b, PP_ALIGN.LEFT, FONT_LIGHT)

# =====================================================================
# SLIDE 33: EVM GOVERNANCE
# =====================================================================
s = make_slide()
add_shape(s, 0, 0, SW, SH, NAVY)
add_shape(s, 0, 0, Inches(0.12), SH, GOLD)
slide_section(s, 'EVM Governance & PMO', 'حوكمة EVM وإدارة مكتب إدارة المشاريع')

gov_items = [
    ('سياسات وإجراءات EVM', 'EVM Policy', NAVY,
     'وضع سياسات واضحة لتطبيق EVM في جميع مشاريع\nالمؤسسة مع تحديد الأدوار والمسؤوليات'),
    ('خط الأساس', 'Baseline Management', TEAL,
     'إدارة خط الأساس للمشروع مع آليات مراقبة\nالتغييرات والتحديثات الرسمية'),
    ('دورة التقارير', 'Reporting Cycle', GOLD,
     'تحديد دورة منتظمة لتقارير EVM (أسبوعي -\nشهري - ربع سنوي) مع قوالب موحدة'),
    ('مراجعة الأداء', 'Performance Reviews', TEAL,
     'اجتماعات دورية لمراجعة مؤشرات الأداء واتخاذ\nالقرارات التصحيحية اللازمة'),
    ('التكامل المؤسسي', 'Enterprise Integration', NAVY,
     'ربط EVM مع الأنظمة المؤسسية ERP، P6، PMIS\nلضمان تدفق البيانات تلقائياً'),
]
for i, (ar, en, color, desc) in enumerate(gov_items):
    x = Inches(0.4) + Inches(i * 2.6)
    add_shape(s, x, Inches(1.6), Inches(2.4), Inches(3.8), color, 0.2)
    add_shape(s, x, Inches(1.6), Inches(2.4), Inches(0.8), color, 0.85)
    add_txt(s, x, Inches(1.65), Inches(2.4), Inches(0.7), ar, 13, WHITE, True, PP_ALIGN.CENTER)
    add_txt(s, x + Inches(0.1), Inches(2.5), Inches(2.2), Inches(0.3), en, 9, RGBColor(220, 220, 220), False, PP_ALIGN.CENTER, FONT_LIGHT)
    add_txt(s, x + Inches(0.1), Inches(3.0), Inches(2.2), Inches(2.0), desc, 9, WHITE, False, PP_ALIGN.CENTER, FONT_LIGHT)

# Bottom governance framework
add_shape(s, Inches(0.4), Inches(5.7), Inches(12.5), Inches(1.3), DARK_NAVY)
add_txt(s, Inches(0.8), Inches(5.8), Inches(5), Inches(0.3), 'PMO Governance Levels', 13, GOLD, True, PP_ALIGN.LEFT)
gov_levels = [
    ('المستوى 1: المشروع', 'Project Level', 'تطبيق EVM على مستوى المشروع الفردي'),
    ('المستوى 2: البرنامج', 'Program Level', 'تجميع مؤشرات EVM لبرنامج متكامل'),
    ('المستوى 3: المحفظة', 'Portfolio Level', 'تقارير موحدة لمحفظة المشاريع'),
]
for i, (ar, en, desc) in enumerate(gov_levels):
    y = Inches(6.1) + Inches(i * 0.3)
    add_txt(s, Inches(0.8), y, Inches(3), Inches(0.28), ar, 10, WHITE, True, PP_ALIGN.RIGHT)
    add_txt(s, Inches(4.0), y, Inches(3), Inches(0.28), f'({en})', 9, TEAL, False, PP_ALIGN.LEFT, FONT_LIGHT)
    add_txt(s, Inches(7.5), y, Inches(5), Inches(0.28), desc, 9, LIGHT_GRAY, False, PP_ALIGN.LEFT, FONT_LIGHT)

# =====================================================================
# SLIDE 34: KEY TAKEAWAYS
# =====================================================================
s = make_slide()
add_full_bg(s, img('02_Art_Park_1920x1080-Gross.jpg'))
add_shape(s, 0, 0, SW, SH, DARK_BG, 0.6)
slide_title_bar(s, 'Key Takeaways & Best Practices', 'الدروس المستفادة وأفضل الممارسات')

add_shape(s, Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.2), NAVY, 0.85)
add_rich(s, Inches(1.2), Inches(1.7), Inches(5.2), Inches(4.8), [
    ('الدروس المستفادة', 18, GOLD, True, FONT),
    ('', 6, WHITE, False, FONT),
    ('✓  ابدأ بـ WBS دقيق ومفصل', 13, WHITE, True, FONT),
    ('     خط الأساس الجيد هو نصف النجاح', 11, LIGHT_GRAY, False, FONT),
    ('', 3, WHITE, False, FONT),
    ('✓  حدد طرق قياس الأداء مسبقاً', 13, WHITE, True, FONT),
    ('     % Complete، Milestones، Units Completed', 11, LIGHT_GRAY, False, FONT),
    ('', 3, WHITE, False, FONT),
    ('✓  استخدم CPI و SPI وليس فقط CV و SV', 13, WHITE, True, FONT),
    ('     المؤشرات النسبية أدق من المطلقة', 11, LIGHT_GRAY, False, FONT),
    ('', 3, WHITE, False, FONT),
    ('✓  حدث البيانات بانتظام', 13, WHITE, True, FONT),
    ('     EVM أداة حية تحتاج تغذية مستمرة', 11, LIGHT_GRAY, False, FONT),
    ('', 3, WHITE, False, FONT),
    ('✓  اربط EVM بإدارة المخاطر', 13, WHITE, True, FONT),
    ('     استخدم التوقعات لتحديد المخاطر مبكراً', 11, LIGHT_GRAY, False, FONT),
])

add_shape(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.2), DARK_NAVY, 0.85)
add_rich(s, Inches(7.2), Inches(1.7), Inches(5.2), Inches(4.8), [
    ('Best Practices', 18, GOLD, True, FONT),
    ('', 6, WHITE, False, FONT),
    ('01  Implement EVM from project initiation', 13, WHITE, True, FONT_LIGHT),
    ('     Don\'t wait until problems appear', 11, LIGHT_GRAY, False, FONT_LIGHT),
    ('', 3, WHITE, False, FONT),
    ('02  Use weighted milestones for measurement', 13, WHITE, True, FONT_LIGHT),
    ('     More accurate than % or 0/100', 11, LIGHT_GRAY, False, FONT_LIGHT),
    ('', 3, WHITE, False, FONT),
    ('03  Integrate EVM with existing systems', 13, WHITE, True, FONT_LIGHT),
    ('     P6, ERP, and PMIS automation', 11, LIGHT_GRAY, False, FONT_LIGHT),
    ('', 3, WHITE, False, FONT),
    ('04  Train your team on EVM concepts', 13, WHITE, True, FONT_LIGHT),
    ('     EVM literacy across all levels', 11, LIGHT_GRAY, False, FONT_LIGHT),
    ('', 3, WHITE, False, FONT),
    ('05  Review and forecast monthly', 13, WHITE, True, FONT_LIGHT),
    ('     Regular EVM reviews with stakeholders', 11, LIGHT_GRAY, False, FONT_LIGHT),
])

# =====================================================================
# SLIDE 35: CLOSING
# =====================================================================
s = make_slide()
add_full_bg(s, img('0000_Overlook_View_1920x1080.jpg'))
add_shape(s, 0, 0, SW, SH, DARK_BG, 0.55)
add_shape(s, 0, 0, Inches(0.12), SH, GOLD)

add_gold_accent(s, Inches(5.5), Inches(2.3), Inches(2.3), Pt(5))

add_rich(s, Inches(1.5), Inches(2.6), Inches(10.3), Inches(3.5), [
    ('إدارة القيمة المكتسبة', 40, WHITE, True, FONT_LIGHT),
    ('Earned Value Management', 18, GOLD, False, FONT_LIGHT),
    ('', 16, WHITE, False, FONT_LIGHT),
    ('في المشاريع الإنشائية', 34, WHITE, True, FONT_LIGHT),
    ('', 20, WHITE, False, FONT_LIGHT),
    ('التحكم في التكاليف  •  إدارة الجدول  •  قياس الأداء  •  التنبؤ بالنتائج', 14, LIGHT_GRAY, False, FONT_LIGHT),
], PP_ALIGN.CENTER)

add_shape(s, Inches(3), Inches(6.0), Inches(7.3), Inches(0.5), NAVY, 0.8)
add_txt(s, Inches(3.2), Inches(6.02), Inches(6.9), Inches(0.45),
    '"ما لا يمكن قياسه، لا يمكن إدارته"  |  PMI Best Practice',
    14, GOLD, True, PP_ALIGN.CENTER, FONT_LIGHT)

# =====================================================================
# SAVE
# =====================================================================
prs.save(OUTPUT)
print(f"Presentation saved to: {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
