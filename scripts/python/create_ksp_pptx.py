from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

IMG_DIR = r"E:\N8N\scraper\scraper2\csi-ultimate\ksp_images"
OUTPUT = r"E:\N8N\scraper\scraper2\csi-ultimate\King_Salman_Park.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SW = prs.slide_width
SH = prs.slide_height

# Colors
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
GOLD = RGBColor(212, 175, 55)
DARK_GOLD = RGBColor(180, 140, 40)
DARK_OVERLAY = RGBColor(10, 20, 15)
TRANSPARENT_BLACK = RGBColor(0, 0, 0)  # used with alpha in overlays

def add_full_bleed_bg(slide, image_path, brightness=1.0):
    slide.shapes.add_picture(image_path, 0, 0, SW, SH)

def add_rect(slide, left, top, width, height, fill_color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if alpha is not None:
        from lxml import etree
        from pptx.oxml.ns import qn
        spPr = shape._element.spPr
        sf = spPr.find(qn('a:solidFill'))
        if sf is not None:
            srgb = sf.find(qn('a:srgbClr'))
            if srgb is not None:
                srgb.set('alpha', str(int(alpha * 1000)))
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri Light'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_rich_text(slide, left, top, width, height, lines, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size, color, bold, font_name) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(6)
    return txBox

def img_path(name):
    return os.path.join(IMG_DIR, name)

# ──────────────────────────────────────────────
# SLIDE 1: COVER
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_full_bleed_bg(slide, img_path("ksp_aerial.jpg"))

# Dark gradient overlay
add_rect(slide, 0, 0, SW, SH, BLACK, alpha=0.55)

# Accent line
add_rect(slide, Inches(1.5), Inches(2.5), Inches(1.2), Pt(4), GOLD)

add_rich_text(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(3), [
    ("King Salman Park", 52, WHITE, True, 'Calibri Light'),
    ("Riyadh, Saudi Arabia", 24, RGBColor(200, 200, 200), False, 'Calibri Light'),
    ("The World's Largest Urban Park", 18, GOLD, False, 'Calibri Light'),
    ("", 12, WHITE, False, 'Calibri Light'),
    ("A Vision of the Saudi Green Initiative  |  Master Plan by Omrania & Henning Larsen", 14, RGBColor(180, 180, 180), False, 'Calibri Light'),
], alignment=PP_ALIGN.LEFT)

# ──────────────────────────────────────────────
# SLIDE 2: OVERVIEW
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_full_bleed_bg(slide, img_path("KSP_Overview.jpg"))
add_rect(slide, 0, 0, SW, SH, BLACK, alpha=0.6)

add_rect(slide, Inches(0.8), Inches(1.2), Inches(0.8), Pt(3), GOLD)

add_rich_text(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5), [
    ("Project Overview", 36, WHITE, True, 'Calibri Light'),
    ("", 10, WHITE, False, 'Calibri Light'),
    ("King Salman Park is a landmark giga-project and a cornerstone of Riyadh's transformation into one of the world's most livable cities.", 16, RGBColor(220, 220, 220), False, 'Calibri Light'),
    ("", 10, WHITE, False, 'Calibri Light'),
    ("•   Total Area: 16.6 km² (4x the size of Central Park, NYC)", 15, RGBColor(200, 200, 200), False, 'Calibri Light'),
    ("•   Location: Heart of Riyadh, along Prince Turki bin Abdulaziz Road", 15, RGBColor(200, 200, 200), False, 'Calibri Light'),
    ("•   Completion: Phase I expected by 2027-2028", 15, RGBColor(200, 200, 200), False, 'Calibri Light'),
    ("•   Cost: Over SAR 23 billion (~$6.1 billion)", 15, RGBColor(200, 200, 200), False, 'Calibri Light'),
    ("•   Designed by: Omrania, Henning Larsen, Gerber Architekten", 15, RGBColor(200, 200, 200), False, 'Calibri Light'),
])

# ──────────────────────────────────────────────
# SLIDE 3: VISION
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_full_bleed_bg(slide, img_path("riba_gerber_original_4.jpg"))
add_rect(slide, 0, 0, SW, SH, BLACK, alpha=0.55)

add_rect(slide, Inches(0.8), Inches(1.2), Inches(0.8), Pt(3), GOLD)

add_rich_text(slide, Inches(0.8), Inches(1.5), Inches(7), Inches(5), [
    ("Vision & Ambition", 36, WHITE, True, 'Calibri Light'),
    ("", 10, WHITE, False, 'Calibri Light'),
    ("Part of the Riyadh Green Initiative — aiming to plant 7.5 million trees across the city.", 16, RGBColor(220, 220, 220), False, 'Calibri Light'),
    ("", 8, WHITE, False, 'Calibri Light'),
    ("The park will feature:", 16, WHITE, True, 'Calibri Light'),
    ("", 6, WHITE, False, 'Calibri Light'),
    ("🌳  A sprawling green oasis with over 200,000 trees", 15, RGBColor(200, 200, 200), False, 'Calibri Light'),
    ("🎨  A world-class cultural district with museums and galleries", 15, RGBColor(200, 200, 200), False, 'Calibri Light'),
    ("🏛️  The Royal Arts Complex & Museum of the Earth", 15, RGBColor(200, 200, 200), False, 'Calibri Light'),
    ("🚶  30 km of pedestrian pathways and cycling tracks", 15, RGBColor(200, 200, 200), False, 'Calibri Light'),
    ("💧  A 1.2 km wadi (valley) with cascading waterfalls", 15, RGBColor(200, 200, 200), False, 'Calibri Light'),
])

# ──────────────────────────────────────────────
# SLIDE 4: MASTERPLAN
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_full_bleed_bg(slide, img_path("ksp_masterplan.jpg"))
add_rect(slide, 0, Inches(5.5), SW, Inches(2), BLACK, alpha=0.7)

add_text_box(slide, Inches(0.8), Inches(5.7), Inches(11), Inches(1.5),
    "Master Plan  —  A seamless integration of nature, culture, and urban living across 16.6 km²",
    18, WHITE, False, PP_ALIGN.LEFT)

# ──────────────────────────────────────────────
# SLIDE 5: WADI & WATERFALL
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_full_bleed_bg(slide, img_path("01_Wadi_Cascade_Waterfall_1920x1080-Gross.jpg"))
add_rect(slide, 0, 0, SW, SH, BLACK, alpha=0.5)

add_rect(slide, Inches(0.8), Inches(1.2), Inches(0.8), Pt(3), GOLD)

add_rich_text(slide, Inches(0.8), Inches(1.5), Inches(7), Inches(4), [
    ("The Wadi & Cascading Waterfall", 36, WHITE, True, 'Calibri Light'),
    ("", 10, WHITE, False, 'Calibri Light'),
    ("A 1.2 km dry valley transformed into a dynamic water feature with terraced waterfalls, lush planting, and pedestrian bridges. The centerpiece of the park's sustainable water management system.", 16, RGBColor(220, 220, 220), False, 'Calibri Light'),
    ("", 10, WHITE, False, 'Calibri Light'),
    ("•   Collects and recycles rainwater", 14, RGBColor(200, 200, 200), False, 'Calibri Light'),
    ("•   Creates a microclimate for surrounding areas", 14, RGBColor(200, 200, 200), False, 'Calibri Light'),
    ("•   Features shaded walkways and viewing platforms", 14, RGBColor(200, 200, 200), False, 'Calibri Light'),
])

# ──────────────────────────────────────────────
# SLIDE 6: ART PARK
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_full_bleed_bg(slide, img_path("02_Art_Park_1920x1080-Gross.jpg"))
add_rect(slide, 0, 0, SW, SH, BLACK, alpha=0.5)

add_rect(slide, Inches(0.8), Inches(1.2), Inches(0.8), Pt(3), GOLD)

add_rich_text(slide, Inches(0.8), Inches(1.5), Inches(7), Inches(4), [
    ("The Art Park", 36, WHITE, True, 'Calibri Light'),
    ("", 10, WHITE, False, 'Calibri Light'),
    ("An open-air museum blending landscape architecture with contemporary sculpture, installations, and interactive art spaces.", 16, RGBColor(220, 220, 220), False, 'Calibri Light'),
    ("", 10, WHITE, False, 'Calibri Light'),
    ("•   Permanent and rotating exhibitions", 14, RGBColor(200, 200, 200), False, 'Calibri Light'),
    ("•   International and local artists", 14, RGBColor(200, 200, 200), False, 'Calibri Light'),
    ("•   Integrated with the park's natural topography", 14, RGBColor(200, 200, 200), False, 'Calibri Light'),
])

# ──────────────────────────────────────────────
# SLIDE 7: ART PARK AERIAL
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_full_bleed_bg(slide, img_path("07_Art_Park_Aerial_View_1920x1080-Gross.jpg"))
add_rect(slide, 0, Inches(5.5), SW, Inches(2), BLACK, alpha=0.7)

add_text_box(slide, Inches(0.8), Inches(5.7), Inches(11), Inches(1.5),
    "Art Park Aerial View  —  Sculpture gardens and cultural pavilions woven into the landscape",
    18, WHITE, False, PP_ALIGN.LEFT)

# ──────────────────────────────────────────────
# SLIDE 8: MUSEUM OF THE EARTH
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_full_bleed_bg(slide, img_path("02_Museum_of_the_Earth_1920x1080-Gross.jpg"))
add_rect(slide, 0, 0, SW, SH, BLACK, alpha=0.5)

add_rect(slide, Inches(0.8), Inches(1.2), Inches(0.8), Pt(3), GOLD)

add_rich_text(slide, Inches(0.8), Inches(1.5), Inches(7), Inches(4), [
    ("Museum of the Earth", 36, WHITE, True, 'Calibri Light'),
    ("", 10, WHITE, False, 'Calibri Light'),
    ("A landmark institution dedicated to Earth sciences, geology, and environmental awareness, designed as an iconic architectural statement within the park.", 16, RGBColor(220, 220, 220), False, 'Calibri Light'),
    ("", 10, WHITE, False, 'Calibri Light'),
    ("•   Immersive exhibits on Earth's geological history", 14, RGBColor(200, 200, 200), False, 'Calibri Light'),
    ("•   Sustainable architecture integrated with the landscape", 14, RGBColor(200, 200, 200), False, 'Calibri Light'),
    ("•   Educational center for climate and environmental research", 14, RGBColor(200, 200, 200), False, 'Calibri Light'),
])

# ──────────────────────────────────────────────
# SLIDE 9: SKYLINE LOOP BRIDGE
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_full_bleed_bg(slide, img_path("02_Skyline_Loop_Bridge_1920x1080-Gross.jpg"))
add_rect(slide, 0, 0, SW, SH, BLACK, alpha=0.5)

add_rect(slide, Inches(0.8), Inches(1.2), Inches(0.8), Pt(3), GOLD)

add_rich_text(slide, Inches(0.8), Inches(1.5), Inches(7), Inches(4), [
    ("Skyline Loop Bridge", 36, WHITE, True, 'Calibri Light'),
    ("", 10, WHITE, False, 'Calibri Light'),
    ("A signature elevated pedestrian bridge offering panoramic views of the Riyadh skyline and the park below. A landmark in its own right.", 16, RGBColor(220, 220, 220), False, 'Calibri Light'),
    ("", 10, WHITE, False, 'Calibri Light'),
    ("•   Continuous loop design — a symbol of unity and connection", 14, RGBColor(200, 200, 200), False, 'Calibri Light'),
    ("•   Accessible for all visitors", 14, RGBColor(200, 200, 200), False, 'Calibri Light'),
    ("•   Illuminated at night as a visual landmark", 14, RGBColor(200, 200, 200), False, 'Calibri Light'),
])

# ──────────────────────────────────────────────
# SLIDE 10: GARDENS & NURSERY
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_full_bleed_bg(slide, img_path("0000_Visitor_Nursery_001_1920x1080.jpg"))
add_rect(slide, 0, 0, SW, SH, BLACK, alpha=0.5)

add_rect(slide, Inches(0.8), Inches(1.2), Inches(0.8), Pt(3), GOLD)

add_rich_text(slide, Inches(0.8), Inches(1.5), Inches(7), Inches(4), [
    ("Thematic Gardens & Nursery", 36, WHITE, True, 'Calibri Light'),
    ("", 10, WHITE, False, 'Calibri Light'),
    ("Diverse botanical gardens showcasing native and adapted plant species, with a dedicated nursery cultivating over 200,000 trees for the park and wider Riyadh greening initiatives.", 16, RGBColor(220, 220, 220), False, 'Calibri Light'),
    ("", 10, WHITE, False, 'Calibri Light'),
    ("•   Native desert flora and drought-resistant species", 14, RGBColor(200, 200, 200), False, 'Calibri Light'),
    ("•   Educational programs on sustainable horticulture", 14, RGBColor(200, 200, 200), False, 'Calibri Light'),
    ("•   A living laboratory for urban greening in arid climates", 14, RGBColor(200, 200, 200), False, 'Calibri Light'),
])

# ──────────────────────────────────────────────
# SLIDE 11: OVERLOOK
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_full_bleed_bg(slide, img_path("0000_Overlook_View_1920x1080.jpg"))
add_rect(slide, 0, Inches(5.5), SW, Inches(2), BLACK, alpha=0.7)

add_text_box(slide, Inches(0.8), Inches(5.7), Inches(11), Inches(1.5),
    "Overlook View  —  Vantage points designed to frame the perfect intersection of nature and city",
    18, WHITE, False, PP_ALIGN.LEFT)

# ──────────────────────────────────────────────
# SLIDE 12: FULL PARK
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_full_bleed_bg(slide, img_path("1509_KSP_Full_Park_001_1920x1080.jpg"))
add_rect(slide, 0, Inches(5.5), SW, Inches(2), BLACK, alpha=0.7)

add_text_box(slide, Inches(0.8), Inches(5.7), Inches(11), Inches(1.5),
    "A New Green Heart for Riyadh  —  16.6 km² of interconnected parks, plazas, and cultural venues",
    18, WHITE, False, PP_ALIGN.LEFT)

# ──────────────────────────────────────────────
# SLIDE 13: GERBER ARCHITEKTEN - DETAIL
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_full_bleed_bg(slide, img_path("riba_gerber_original_4.jpg"))
add_rect(slide, 0, 0, SW, SH, BLACK, alpha=0.5)

add_rect(slide, Inches(0.8), Inches(1.2), Inches(0.8), Pt(3), GOLD)

add_rich_text(slide, Inches(0.8), Inches(1.5), Inches(7), Inches(4), [
    ("Architectural Excellence", 36, WHITE, True, 'Calibri Light'),
    ("", 10, WHITE, False, 'Calibri Light'),
    ("Gerber Architekten's contribution to King Salman Park: innovative design language that bridges modern architecture with the Saudi cultural identity.", 16, RGBColor(220, 220, 220), False, 'Calibri Light'),
    ("", 10, WHITE, False, 'Calibri Light'),
    ("•   Fluid, organic forms inspired by natural desert landscapes", 14, RGBColor(200, 200, 200), False, 'Calibri Light'),
    ("•   Cutting-edge sustainable building technologies", 14, RGBColor(200, 200, 200), False, 'Calibri Light'),
    ("•   Integration with the park's topography and microclimates", 14, RGBColor(200, 200, 200), False, 'Calibri Light'),
])

# ──────────────────────────────────────────────
# SLIDE 14: SUSTAINABILITY
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_full_bleed_bg(slide, img_path("omrania_ksp_render_2.jpg"))
add_rect(slide, 0, 0, SW, SH, BLACK, alpha=0.55)

add_rect(slide, Inches(0.8), Inches(1.2), Inches(0.8), Pt(3), GOLD)

add_rich_text(slide, Inches(0.8), Inches(1.5), Inches(7), Inches(5), [
    ("Sustainability at Scale", 36, WHITE, True, 'Calibri Light'),
    ("", 10, WHITE, False, 'Calibri Light'),
    ("King Salman Park is a model for sustainable urban development in arid regions.", 16, RGBColor(220, 220, 220), False, 'Calibri Light'),
    ("", 8, WHITE, False, 'Calibri Light'),
    ("♻️  Water Conservation: Greywater recycling and smart irrigation", 15, RGBColor(200, 200, 200), False, 'Calibri Light'),
    ("🌱  Carbon Sink: 200,000+ trees absorbing CO₂ and cooling the city", 15, RGBColor(200, 200, 200), False, 'Calibri Light'),
    ("☀️  Renewable Energy: Solar-powered facilities and lighting", 15, RGBColor(200, 200, 200), False, 'Calibri Light'),
    ("🚲  Mobility: Car-free zones, 30 km of pedestrian/cycle paths", 15, RGBColor(200, 200, 200), False, 'Calibri Light'),
    ("🏗️  Local Materials: Saudi-sourced stone and construction materials", 15, RGBColor(200, 200, 200), False, 'Calibri Light'),
])

# ──────────────────────────────────────────────
# SLIDE 15: CLOSING
# ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_full_bleed_bg(slide, img_path("1509_KSP_Park_Overview_001_1920x1080.jpg"))
add_rect(slide, 0, 0, SW, SH, BLACK, alpha=0.6)

add_rect(slide, Inches(5.5), Inches(2.8), Inches(2.3), Pt(4), GOLD)

add_rich_text(slide, Inches(2), Inches(3.1), Inches(9.3), Inches(3), [
    ("King Salman Park", 44, WHITE, True, 'Calibri Light'),
    ("Where Nature Meets Tomorrow", 22, GOLD, False, 'Calibri Light'),
    ("", 16, WHITE, False, 'Calibri Light'),
    ("Riyadh's Green Heart  |  A Saudi Vision 2030 Legacy Project", 16, RGBColor(180, 180, 180), False, 'Calibri Light'),
], alignment=PP_ALIGN.CENTER)

prs.save(OUTPUT)
print(f"Presentation saved to: {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
